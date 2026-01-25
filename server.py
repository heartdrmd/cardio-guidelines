#!/usr/bin/env python3
"""
Cardiology Guidelines AI Server

Flask API that:
1. Accepts PDF uploads (single or multiple)
2. Extracts text and chunks it
3. Generates embeddings and stores in Postgres (pgvector)
4. Answers questions using Claude + relevant guideline context
5. Tracks all PDFs and allows deletion/replacement
"""

import os
import re
import hashlib
from datetime import datetime
from flask import Flask, jsonify, request, send_from_directory
from flask_cors import CORS
import psycopg
from psycopg.rows import dict_row
import pdfplumber
import anthropic

# Test pptx import at startup
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor  # Note: RGBColor not RgbColor
    PPTX_AVAILABLE = True
    print("✓ python-pptx loaded successfully")
except ImportError as e:
    PPTX_AVAILABLE = False
    print(f"✗ python-pptx import failed: {e}")

# =============================================================================
# CONFIGURATION
# =============================================================================

app = Flask(__name__)
CORS(app)

# Database
DATABASE_URL = os.environ.get('DATABASE_URL')

# Anthropic
ANTHROPIC_API_KEY = os.environ.get('ANTHROPIC_API_KEY')
claude = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY) if ANTHROPIC_API_KEY else None

# Upload folder
UPLOAD_FOLDER = '/tmp/guidelines'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Chunking config
CHUNK_SIZE = 1500  # characters per chunk
CHUNK_OVERLAP = 200  # overlap between chunks

# =============================================================================
# TEXT SPACING FIX (for poorly extracted PDF text)
# =============================================================================

_SPACING_WORDS = set('''
a about above according additionally after age all also alternatives among an and any 
appropriate are as assessment associated at available based be because been before being 
benefit between both but by can cardiac choice class clinical combined complication 
complications condition conditions consider considered continue controlled coronary could 
current data day days decision degree despite determined develop diagnostic disease 
diseases do does dose dosing drug during each effective elderly elevated especially 
established evaluate evaluated evaluating event evidence exclude factor factors failure 
fibrillation first follow following for from general generally give given good greater 
guideline guidelines harm have having heart high higher history however i if important 
improve in include including increase increased indication indications individual 
information initial intervention interventions into is it its known least less level 
levels life likely limited long low lower made major make making management may mean 
mechanical medical medication moderate months more mortality most my necessary need new 
no normal not now number observation observed of often older on one only option options 
or other outcome outcomes out over overall part particular patient patients per percent 
performed performing place population possible possibly potential present previous prior 
probably procedure procedures prognosis prosthetic provide provided provides quality 
randomized rate rather reasonable recommend recommendation recommendations recommended 
recommending reduce reduced regurgitation related repair replacement require required 
requires requiring result results review risk risks safe safety same score select 
selected selection serial setting severe severity several she shared sharing should 
show shown similar since so some specific stenosis strength strong studies study 
subsequent such suggest suggested suggests support supported surgery surgical symptoms 
systematic than that the therapeutic therapy their them then there therefore these they 
this those three through time to together total toward treatment trial trials two 
undergoing unless until up upon us use used using usually valve valves value values 
various versus very was weak well were what when where whether which while who will 
with within without women would year years yet you young your unchanged
anticoagulant anticoagulation anticoagulants aortic atrial bioprosthetic bypass 
comorbidities contraindicated contraindication durability ejection embolism expectancy
fraction individualize individualized intravenous ischemia ischemic 
leaflet leaflets lifelong managed medically mitral native noncardiac obstructive
percutaneous pharmacological postoperative preoperative pulmonary reintervention
reoperation stent stented stents stroke structural subvalvular thromboembolism thrombosis
thromboembolic transcatheter tricuspid valvular warfarin bioprosthesis
accounts appropriately desires discussion discussions includes preferences type types
process processes criteria criterion presence absence remain remains remaining choose
inpatients left ventricular right atrium ventricle arrhythmia arrhythmias rhythm rhythms
'''.lower().split())

def _segment_text_dp(text):
    """Use DP to segment lowercase text into words."""
    text_lower = text.lower()
    n = len(text)
    if n == 0:
        return ""
    INF = float('inf')
    dp = [(INF, -1)] * (n + 1)
    dp[0] = (0, -1)
    for i in range(1, n + 1):
        for j in range(max(0, i - 25), i):
            word = text_lower[j:i]
            word_len = i - j
            if word in _SPACING_WORDS:
                cost = -word_len * 10
            elif word_len == 1 and word in 'ai':
                cost = 0
            elif word_len == 1:
                cost = 5
            else:
                cost = word_len * 3
            new_cost = dp[j][0] + cost
            if new_cost < dp[i][0]:
                dp[i] = (new_cost, j)
    segments = []
    i = n
    while i > 0:
        prev = dp[i][1]
        segments.append(text[prev:i])
        i = prev
    segments.reverse()
    result = []
    for seg in segments:
        if len(seg) == 1 and seg.lower() not in 'ai' and seg.lower() not in _SPACING_WORDS:
            if result:
                result[-1] += seg
            else:
                result.append(seg)
        else:
            result.append(seg)
    return ' '.join(result)

def _process_spacing_token(token):
    """Process a single token for spacing fix."""
    if len(token) < 8:
        return token
    leading, trailing, core = '', '', token
    match = re.match(r'^([^a-zA-Z]*)', token)
    if match:
        leading = match.group(1)
        core = token[len(leading):]
    match = re.search(r'([^a-zA-Z]*)$', core)
    if match:
        trailing = match.group(1)
        core = core[:len(core)-len(trailing)]
    if len(core) < 8:
        return token
    if re.search(r'[<>=]', core):
        parts = re.split(r'([<>=]+)', core)
        processed = []
        for part in parts:
            if re.match(r'^[<>=]+$', part):
                processed.append(part)
            elif len(part) >= 8:
                processed.append(_process_spacing_token(part))
            else:
                processed.append(part)
        return leading + ''.join(processed) + trailing
    if re.match(r'^[a-z]+$', core):
        return leading + _segment_text_dp(core) + trailing
    if re.match(r'^[A-Z][a-z]{7,}$', core):
        segmented = _segment_text_dp(core.lower())
        if segmented:
            segmented = segmented[0].upper() + segmented[1:]
        return leading + segmented + trailing
    if re.search(r'[a-z][A-Z]', core):
        parts = re.split(r'(?<=[a-z])(?=[A-Z])', core)
        processed = []
        for part in parts:
            if len(part) >= 8 and re.match(r'^[A-Z]?[a-z]+$', part):
                seg = _segment_text_dp(part.lower())
                if part[0].isupper() and seg:
                    seg = seg[0].upper() + seg[1:]
                processed.append(seg)
            else:
                processed.append(part)
        return leading + ' '.join(processed) + trailing
    return token

def fix_text_spacing(text):
    """Fix missing spaces in poorly extracted PDF text using DP algorithm."""
    if not text:
        return text
    text = re.sub(r'([,;:.!?])([A-Za-z])', r'\1 \2', text)
    tokens = text.split()
    result = []
    for token in tokens:
        if "'" in token:
            parts = token.split("'")
            fixed = [_process_spacing_token(p) for p in parts]
            result.append("'".join(fixed))
        else:
            result.append(_process_spacing_token(token))
    text = ' '.join(result)
    text = re.sub(r'  +', ' ', text)
    return text.strip()

def fix_text_spacing_with_ai(text, use_cache=True):
    """
    Fix missing spaces using AI - handles ANY vocabulary perfectly.
    More expensive but 100% accurate.
    """
    if not text or len(text) < 20:
        return text
    
    # Quick check - does it even need fixing?
    words = text.split()
    long_words = [w for w in words if len(w) > 35 and re.search(r'[a-z]{15,}', w)]
    if not long_words:
        return text  # No obvious issues, skip AI call
    
    if not claude:
        # Fall back to DP algorithm if no API key
        return fix_text_spacing(text)
    
    try:
        response = claude.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=len(text) + 500,
            messages=[{
                "role": "user", 
                "content": f"""Fix the spacing in this text extracted from a PDF. Words are concatenated without spaces.

RULES:
1. Add spaces where words run together
2. Keep all original words and meaning
3. Preserve punctuation and numbers
4. Return ONLY the fixed text

TEXT: {text}

FIXED:"""
            }]
        )
        return response.content[0].text.strip()
    except Exception as e:
        print(f"AI spacing fix failed: {e}, using DP fallback")
        return fix_text_spacing(text)

def check_text_needs_spacing_fix(text):
    """Check if text has spacing issues (long concatenated words)."""
    if not text:
        return False
    words = text.split()
    long_words = [w for w in words if len(w) > 35 and re.search(r'[a-z]{15,}', w)]
    return len(long_words) > 0

# =============================================================================
# DATABASE HELPERS
# =============================================================================

def get_db():
    """Get database connection."""
    conn = psycopg.connect(DATABASE_URL)
    return conn

def init_db():
    """Initialize database tables if they don't exist."""
    conn = get_db()
    cur = conn.cursor()
    
    # Guidelines table - tracks each uploaded PDF
    cur.execute("""
        CREATE TABLE IF NOT EXISTS guidelines (
            id SERIAL PRIMARY KEY,
            name TEXT NOT NULL,
            year INT,
            source TEXT,
            filename TEXT,
            file_hash TEXT UNIQUE,
            page_count INT,
            extracted_title TEXT,
            extracted_headers TEXT,
            file_size_bytes INT,
            pdf_data BYTEA,
            created_at TIMESTAMP DEFAULT NOW(),
            updated_at TIMESTAMP DEFAULT NOW()
        );
    """)
    
    # Chunks table - stores extracted text with embeddings
    cur.execute("""
        CREATE TABLE IF NOT EXISTS chunks (
            id SERIAL PRIMARY KEY,
            guideline_id INT REFERENCES guidelines(id) ON DELETE CASCADE,
            section TEXT,
            page_number INT,
            content TEXT NOT NULL,
            embedding vector(1536),
            is_reference BOOLEAN DEFAULT FALSE,
            created_at TIMESTAMP DEFAULT NOW()
        );
    """)
    
    # Create index if not exists
    cur.execute("""
        CREATE INDEX IF NOT EXISTS chunks_embedding_idx 
        ON chunks USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100);
    """)
    
    # Index for fast guideline lookups
    cur.execute("""
        CREATE INDEX IF NOT EXISTS chunks_guideline_idx ON chunks(guideline_id);
    """)
    
    # Topics table - stores AI-extracted topics for each guideline
    cur.execute("""
        CREATE TABLE IF NOT EXISTS guideline_topics (
            id SERIAL PRIMARY KEY,
            guideline_id INT REFERENCES guidelines(id) ON DELETE CASCADE,
            topic_name TEXT NOT NULL,
            icon TEXT DEFAULT '📋',
            start_page INT,
            end_page INT,
            chunk_ids TEXT,
            created_at TIMESTAMP DEFAULT NOW()
        );
    """)
    
    # PPTX Jobs table - for background processing
    cur.execute("""
        CREATE TABLE IF NOT EXISTS pptx_jobs (
            id SERIAL PRIMARY KEY,
            guideline_id INT REFERENCES guidelines(id) ON DELETE CASCADE,
            topic TEXT,
            model TEXT,
            status TEXT DEFAULT 'pending',
            progress INT DEFAULT 0,
            total_batches INT DEFAULT 0,
            result_data TEXT,
            error TEXT,
            cost FLOAT,
            tokens INT,
            created_at TIMESTAMP DEFAULT NOW(),
            completed_at TIMESTAMP
        );
    """)
    
    # Spacing fix jobs table - for background processing of chunk fixes
    cur.execute("""
        CREATE TABLE IF NOT EXISTS spacing_jobs (
            id SERIAL PRIMARY KEY,
            guideline_id INT REFERENCES guidelines(id) ON DELETE CASCADE,
            status TEXT DEFAULT 'pending',
            progress INT DEFAULT 0,
            total_chunks INT DEFAULT 0,
            fixed_chunks INT DEFAULT 0,
            error TEXT,
            created_at TIMESTAMP DEFAULT NOW(),
            completed_at TIMESTAMP
        );
    """)
    
    # Migration: Add new columns if they don't exist
    # Add pdf_data column to guidelines
    try:
        cur.execute("ALTER TABLE guidelines ADD COLUMN IF NOT EXISTS pdf_data BYTEA")
    except:
        pass  # Column may already exist
    
    # Add is_reference column to chunks
    try:
        cur.execute("ALTER TABLE chunks ADD COLUMN IF NOT EXISTS is_reference BOOLEAN DEFAULT FALSE")
    except:
        pass  # Column may already exist
    
    # Add pptx_data and html_data columns to pptx_jobs for caching generated files
    try:
        cur.execute("ALTER TABLE pptx_jobs ADD COLUMN IF NOT EXISTS pptx_data BYTEA")
    except:
        pass
    
    try:
        cur.execute("ALTER TABLE pptx_jobs ADD COLUMN IF NOT EXISTS html_data TEXT")
    except:
        pass
    
    conn.commit()
    cur.close()
    conn.close()

# =============================================================================
# PDF PROCESSING
# =============================================================================

def fix_pdf_spacing(text):
    """Fix missing spaces in PDF-extracted text.
    Uses DP algorithm for light fixes, AI for severe cases."""
    if not text:
        return text
    
    # Check if it needs heavy-duty AI fixing
    if check_text_needs_spacing_fix(text):
        # Use AI for severe cases (many long concatenated words)
        return fix_text_spacing_with_ai(text)
    else:
        # Use fast DP algorithm for minor issues
        return fix_text_spacing(text)

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF with page numbers and headers."""
    pages = []
    all_headers = []
    title = None
    
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            
            # Fix spacing issues common in PDF extraction
            text = fix_pdf_spacing(text)
            
            if text.strip():
                pages.append({
                    'page_number': i + 1,
                    'text': text
                })
                
                # Try to extract title from first page
                if i == 0 and not title:
                    lines = text.strip().split('\n')
                    if lines:
                        # First non-empty line is often the title
                        title = lines[0][:200]  # Limit length
                
                # Extract potential headers (lines that look like headings)
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    # Heuristic: short lines in caps or with numbers might be headers
                    if line and len(line) < 100:
                        if line.isupper() or re.match(r'^\d+\.?\s+[A-Z]', line):
                            all_headers.append(line)
    
    return pages, title, all_headers[:50]  # Limit headers

def chunk_text(text, page_number, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """Split text into overlapping chunks."""
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        
        # Try to break at sentence boundary
        if end < len(text):
            last_period = chunk.rfind('.')
            last_newline = chunk.rfind('\n')
            break_point = max(last_period, last_newline)
            if break_point > chunk_size * 0.5:
                chunk = chunk[:break_point + 1]
                end = start + break_point + 1
        
        if chunk.strip():
            chunks.append({
                'content': chunk.strip(),
                'page_number': page_number
            })
        
        start = end - overlap
        if start < 0:
            start = 0
        if end >= len(text):
            break
    
    return chunks

def get_embedding(text):
    """Generate embedding for a single text using OpenAI API."""
    import requests
    
    openai_key = os.environ.get('OPENAI_API_KEY')
    if not openai_key:
        return None
    
    try:
        response = requests.post(
            'https://api.openai.com/v1/embeddings',
            headers={
                'Authorization': f'Bearer {openai_key}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'text-embedding-3-small',
                'input': text[:8000]
            },
            timeout=30
        )
        
        if response.status_code == 200:
            return response.json()['data'][0]['embedding']
    except Exception as e:
        print(f"Embedding error: {e}")
    
    return None

def get_embeddings_batch(texts):
    """Generate embeddings for multiple texts in one API call (up to 2048)."""
    import requests
    
    openai_key = os.environ.get('OPENAI_API_KEY')
    if not openai_key:
        return [None] * len(texts)
    
    # Truncate each text to 8000 chars
    truncated = [t[:8000] for t in texts]
    
    try:
        response = requests.post(
            'https://api.openai.com/v1/embeddings',
            headers={
                'Authorization': f'Bearer {openai_key}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'text-embedding-3-small',
                'input': truncated
            },
            timeout=120  # Longer timeout for batch
        )
        
        if response.status_code == 200:
            data = response.json()['data']
            # Sort by index to maintain order
            sorted_data = sorted(data, key=lambda x: x['index'])
            return [item['embedding'] for item in sorted_data]
        else:
            print(f"Batch embedding error: {response.status_code} - {response.text}")
    except Exception as e:
        print(f"Batch embedding error: {e}")
    
    return [None] * len(texts)

def ingest_pdf(pdf_path, name, year=None, source=None):
    """Process PDF and store chunks in database."""
    conn = get_db()
    cur = conn.cursor()
    
    # Get file info
    file_size = os.path.getsize(pdf_path)
    
    # Read PDF binary data
    with open(pdf_path, 'rb') as f:
        pdf_binary = f.read()
    
    # Calculate file hash to prevent duplicates
    file_hash = hashlib.md5(pdf_binary).hexdigest()
    
    # Check if already ingested
    cur.execute("SELECT id, name FROM guidelines WHERE file_hash = %s", (file_hash,))
    existing = cur.fetchone()
    if existing:
        cur.close()
        conn.close()
        return {'status': 'exists', 'guideline_id': existing[0], 'name': existing[1]}
    
    # Extract text and metadata
    pages, extracted_title, headers = extract_text_from_pdf(pdf_path)
    if not pages:
        cur.close()
        conn.close()
        return {'status': 'error', 'message': 'No text extracted from PDF'}
    
    # Use extracted title if no name provided
    if not name or name == os.path.basename(pdf_path).replace('.pdf', ''):
        if extracted_title:
            name = extracted_title
    
    # Insert guideline record WITH PDF data
    cur.execute("""
        INSERT INTO guidelines (name, year, source, filename, file_hash, page_count, 
                                extracted_title, extracted_headers, file_size_bytes, pdf_data)
        VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
        RETURNING id
    """, (name, year, source, os.path.basename(pdf_path), file_hash, 
          len(pages), extracted_title, '\n'.join(headers), file_size, pdf_binary))
    guideline_id = cur.fetchone()[0]
    
    # Detect reference section start page - be MORE specific
    # Only look for actual reference/bibliography section, not author info
    reference_start_page = None
    for page in pages:
        text_lower = page['text'].lower()
        # Only match actual reference sections - must be near end of document
        if page['page_number'] > len(pages) * 0.6:  # Only look in last 40% of document
            # Check for reference section headers
            if any(marker in text_lower for marker in [
                '\nreferences\n', 
                '\n\nreferences\n',
                'references\n1.', 
                'references\n 1.',
                '\nbibliography\n',
                '\nliterature cited\n',
            ]):
                reference_start_page = page['page_number']
                print(f"Detected reference section starting at page {reference_start_page}")
                break
    
    # Helper function to detect if a chunk is likely a reference/citation
    def is_reference_chunk(text):
        """Detect if text chunk contains journal citations rather than clinical content."""
        text_lower = text.lower()
        
        # Count citation patterns
        citation_patterns = 0
        
        # Journal name patterns (common cardiology journals)
        journal_names = ['circulation', 'j am coll cardiol', 'jacc', 'eur heart j', 
                        'heart rhythm', 'n engl j med', 'nejm', 'lancet', 'jama',
                        'am j cardiol', 'heart', 'europace', 'j cardiovasc', 'ann intern med']
        for journal in journal_names:
            if journal in text_lower:
                citation_patterns += text_lower.count(journal)
        
        # Year patterns like "2020;" or "2019:" or "(2020)"
        import re
        year_patterns = len(re.findall(r'(?:19|20)\d{2}[;:\)\.]', text))
        citation_patterns += year_patterns
        
        # DOI patterns
        doi_count = text_lower.count('doi:') + text_lower.count('doi.org') + text_lower.count('10.1')
        citation_patterns += doi_count * 2  # Weight DOIs heavily
        
        # Volume/page patterns like "142:1234" or "75(4):123-456"
        vol_page = len(re.findall(r'\d+[:\(]\d+[\):]?\d*[-–]\d+', text))
        citation_patterns += vol_page
        
        # Author patterns: "Smith J," or "Smith JA," or "et al"
        author_patterns = len(re.findall(r'[A-Z][a-z]+\s[A-Z]{1,2}[,\.]', text))
        author_patterns += text_lower.count('et al')
        citation_patterns += author_patterns // 2  # Don't over-weight
        
        # Numbered reference patterns: "1." "2." at start of lines
        numbered_refs = len(re.findall(r'(?:^|\n)\s*\d{1,3}\.\s*[A-Z]', text))
        citation_patterns += numbered_refs
        
        # If high density of citation patterns, it's likely references
        # Threshold: more than 3 patterns per 500 characters
        density = citation_patterns / (len(text) / 500) if len(text) > 0 else 0
        
        return density > 2.5
    
    # Collect all chunks first
    all_chunks = []
    for page in pages:
        chunks = chunk_text(page['text'], page['page_number'])
        for chunk in chunks:
            # Mark as reference if at or after reference section OR if chunk looks like citations
            is_ref = reference_start_page and page['page_number'] >= reference_start_page
            
            # Also check chunk content for citation patterns (catches inline references)
            if not is_ref and is_reference_chunk(chunk['content']):
                is_ref = True
                
            chunk['is_reference'] = is_ref
        all_chunks.extend(chunks)
    
    # Fix spacing issues in chunks BEFORE saving to database (BATCHED for speed)
    # This ensures clean data is stored, not patched later
    chunks_to_fix = [c for c in all_chunks if check_text_needs_spacing_fix(c['content'])]
    chunks_fixed = 0
    
    if chunks_to_fix:
        print(f"  Found {len(chunks_to_fix)} chunks with spacing issues, fixing with AI...")
        SPACING_BATCH_SIZE = 5
        
        for batch_start in range(0, len(chunks_to_fix), SPACING_BATCH_SIZE):
            batch = chunks_to_fix[batch_start:batch_start + SPACING_BATCH_SIZE]
            
            try:
                if claude:
                    # Build batch prompt
                    batch_texts = []
                    for i, chunk in enumerate(batch):
                        batch_texts.append(f"[CHUNK {i+1}]\n{chunk['content']}\n[/CHUNK {i+1}]")
                    
                    combined_text = "\n\n".join(batch_texts)
                    
                    response = claude.messages.create(
                        model="claude-sonnet-4-20250514",
                        max_tokens=len(combined_text) + 2000,
                        messages=[{
                            "role": "user",
                            "content": f"""Fix the spacing in these text chunks. Words are concatenated without spaces.
Add spaces where words run together. Keep [CHUNK N] markers exactly.

{combined_text}

Return fixed chunks with same markers:"""
                        }]
                    )
                    
                    result = response.content[0].text
                    
                    # Parse and apply fixes
                    for i, chunk in enumerate(batch):
                        pattern = rf'\[CHUNK {i+1}\]\s*(.*?)\s*\[/CHUNK {i+1}\]'
                        match = re.search(pattern, result, re.DOTALL)
                        if match:
                            fixed = match.group(1).strip()
                            if fixed and fixed != chunk['content']:
                                chunk['content'] = fixed
                                chunks_fixed += 1
                else:
                    # No AI - use DP fallback
                    for chunk in batch:
                        fixed = fix_text_spacing(chunk['content'])
                        if fixed != chunk['content']:
                            chunk['content'] = fixed
                            chunks_fixed += 1
                            
            except Exception as e:
                print(f"  Batch AI fix failed: {e}, using DP fallback")
                for chunk in batch:
                    try:
                        fixed = fix_text_spacing(chunk['content'])
                        if fixed != chunk['content']:
                            chunk['content'] = fixed
                            chunks_fixed += 1
                    except:
                        pass
        
        print(f"  Fixed spacing in {chunks_fixed} of {len(chunks_to_fix)} chunks")
    
    # Get embeddings in batches of 500
    BATCH_SIZE = 500
    all_embeddings = []
    
    for i in range(0, len(all_chunks), BATCH_SIZE):
        batch = all_chunks[i:i + BATCH_SIZE]
        texts = [c['content'] for c in batch]
        embeddings = get_embeddings_batch(texts)
        all_embeddings.extend(embeddings)
    
    # Insert all chunks with embeddings
    total_chunks = 0
    ref_chunks = 0
    for chunk, embedding in zip(all_chunks, all_embeddings):
        is_ref = chunk.get('is_reference', False)
        if is_ref:
            ref_chunks += 1
            
        if embedding:
            cur.execute("""
                INSERT INTO chunks (guideline_id, page_number, content, embedding, is_reference)
                VALUES (%s, %s, %s, %s, %s)
            """, (guideline_id, chunk['page_number'], chunk['content'], str(embedding), is_ref))
        else:
            cur.execute("""
                INSERT INTO chunks (guideline_id, page_number, content, is_reference)
                VALUES (%s, %s, %s, %s)
            """, (guideline_id, chunk['page_number'], chunk['content'], is_ref))
        
        total_chunks += 1
    
    conn.commit()
    cur.close()
    conn.close()
    
    print(f"Ingested {total_chunks} chunks ({ref_chunks} marked as references, {chunks_fixed} spacing-fixed)")
    
    return {
        'status': 'success',
        'guideline_id': guideline_id,
        'name': name,
        'pages': len(pages),
        'chunks': total_chunks,
        'reference_chunks': ref_chunks,
        'chunks_spacing_fixed': chunks_fixed,
        'extracted_title': extracted_title
    }

# =============================================================================
# SEARCH
# =============================================================================

def understand_query_with_ai(query, model):
    """Use the selected AI model to truly understand the clinical question and extract search terms."""
    
    prompt = f"""You are a medical search expert. Analyze this cardiology question and extract the BEST search terms.

QUESTION: "{query}"

YOUR TASK:
1. UNDERSTAND what the clinician is really asking
2. EXPAND all medical abbreviations to their full terms:
   - ASA → aspirin
   - AFib/AF → atrial fibrillation
   - HTN → hypertension
   - DM → diabetes mellitus
   - CAD → coronary artery disease
   - MI → myocardial infarction
   - HF/CHF → heart failure
   - DOAC/NOAC → direct oral anticoagulant (also add specific drugs: apixaban, rivaroxaban, dabigatran)
   - ACE-I → ACE inhibitor
   - ARB → angiotensin receptor blocker
   - BB → beta blocker
   - CCB → calcium channel blocker
   - etc.

3. IDENTIFY the core clinical concepts:
   - Specific medications or drug classes
   - Diseases or conditions
   - Procedures or interventions
   - Clinical scores or criteria
   - Patient populations

4. IGNORE generic words that won't help searching:
   - Question words (when, what, how, should, etc.)
   - Generic verbs (use, give, take, treat, manage, etc.)
   - Articles and prepositions

Return ONLY a comma-separated list of search terms. Nothing else.

EXAMPLES:
- "When do I use ASA in afib?" → aspirin, atrial fibrillation, antiplatelet, stroke prevention
- "DOAC vs warfarin for mechanical valve" → direct oral anticoagulant, apixaban, rivaroxaban, warfarin, mechanical valve, prosthetic valve
- "HTN management in pregnancy" → hypertension, pregnancy, antihypertensive, preeclampsia, gestational hypertension
- "When to start BB after MI" → beta blocker, metoprolol, carvedilol, myocardial infarction, post-MI, secondary prevention"""

    try:
        if model.startswith('gpt-'):
            # Use OpenAI
            import requests as req
            openai_key = os.environ.get('OPENAI_API_KEY')
            resp = req.post(
                'https://api.openai.com/v1/chat/completions',
                headers={'Authorization': f'Bearer {openai_key}', 'Content-Type': 'application/json'},
                json={
                    'model': model,
                    'messages': [{'role': 'user', 'content': prompt}],
                    'temperature': 0,
                    'max_completion_tokens': 200
                },
                timeout=30
            )
            if resp.status_code == 200:
                terms = resp.json()['choices'][0]['message']['content'].strip()
            else:
                print(f"OpenAI query understanding failed: {resp.status_code}")
                return None
        else:
            # Use Claude (strip -thinking suffix if present)
            actual_model = model.replace('-thinking', '')
            response = claude.messages.create(
                model=actual_model,
                max_tokens=200,
                temperature=0,
                messages=[{"role": "user", "content": prompt}]
            )
            terms = response.content[0].text.strip()
        
        print(f"AI ({model}) understood query as: {terms}")
        return [t.strip().lower() for t in terms.split(',') if t.strip()]
    
    except Exception as e:
        print(f"Query understanding failed: {e}")
        return None


def search_guidelines(query, limit=50, model="claude-sonnet-4-20250514"):
    """Search for relevant guideline chunks using AI-powered hybrid search."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Step 1: Use the SELECTED AI MODEL to understand the query
    ai_terms = understand_query_with_ai(query, model)
    
    if ai_terms:
        print(f"AI search terms: {ai_terms}")
    else:
        # Fallback: just use the original query words (3+ chars)
        ai_terms = [w.lower().strip('?.,/') for w in query.split() if len(w) >= 3]
        print(f"Fallback search terms: {ai_terms}")
    
    # Step 2: Vector search with original query (semantic similarity)
    embedding = get_embedding(query)
    
    print(f"=== SEARCH DEBUG ===")
    print(f"Query: {query}")
    print(f"Model for understanding: {model}")
    print(f"Search terms: {ai_terms}")
    
    results = []
    seen_ids = set()
    
    if embedding:
        # Vector similarity search - EXCLUDE reference sections
        # Format embedding as string for pgvector
        emb_str = '[' + ','.join(map(str, embedding)) + ']'
        cur.execute("""
            SELECT c.id, c.content, c.page_number, c.section,
                   g.id as guideline_id, g.name as guideline_name, g.year, g.source,
                   1 - (c.embedding <=> %s::vector) as similarity
            FROM chunks c
            JOIN guidelines g ON c.guideline_id = g.id
            WHERE c.embedding IS NOT NULL
            AND (c.is_reference IS NULL OR c.is_reference = FALSE)
            ORDER BY c.embedding <=> %s::vector
            LIMIT %s
        """, (emb_str, emb_str, limit))
        results = cur.fetchall()
        print(f"Vector search returned: {len(results)} results")
        seen_ids = {r['id'] for r in results}
    else:
        print("No embedding - skipping vector search")
    
    # Step 3: Keyword search for each AI-extracted term - EXCLUDE references
    for term in ai_terms[:10]:  # Use top 10 AI terms
        if len(term) >= 3:
            cur.execute("""
                SELECT c.id, c.content, c.page_number, c.section,
                       g.id as guideline_id, g.name as guideline_name, g.year, g.source,
                       0.7 as similarity
                FROM chunks c
                JOIN guidelines g ON c.guideline_id = g.id
                WHERE LOWER(c.content) LIKE %s
                AND (c.is_reference IS NULL OR c.is_reference = FALSE)
                LIMIT %s
            """, (f'%{term}%', limit // 2))
            keyword_results = cur.fetchall()
            
            for r in keyword_results:
                if r['id'] not in seen_ids:
                    results.append(r)
                    seen_ids.add(r['id'])
    
    cur.close()
    conn.close()
    
    # Sort by similarity and return top results
    results.sort(key=lambda x: x.get('similarity', 0), reverse=True)
    return results[:limit]

# =============================================================================
# PRICING (per 1M tokens)
# =============================================================================

MODEL_PRICING = {
    # Claude models (input, output)
    'claude-opus-4-20250514': {'input': 15.00, 'output': 75.00},
    'claude-sonnet-4-20250514': {'input': 3.00, 'output': 15.00},
    'claude-haiku-3-5-20241022': {'input': 0.80, 'output': 4.00},
    # OpenAI models
    'gpt-5.2': {'input': 1.75, 'output': 14.00},
}

def calculate_cost(model, input_tokens, output_tokens):
    """Calculate cost in dollars for API call."""
    pricing = MODEL_PRICING.get(model.replace('-thinking', ''), {'input': 3.00, 'output': 15.00})
    input_cost = (input_tokens / 1_000_000) * pricing['input']
    output_cost = (output_tokens / 1_000_000) * pricing['output']
    return round(input_cost + output_cost, 6)

# =============================================================================
# CLAUDE INTEGRATION
# =============================================================================

def ask_claude(question, context_chunks, model="claude-sonnet-4-20250514"):
    """Ask Claude or GPT a question with guideline context."""
    
    # Build context from chunks
    context_parts = []
    for chunk in context_chunks:
        source = f"{chunk['guideline_name']}"
        if chunk.get('year'):
            source += f" ({chunk['year']})"
        if chunk.get('source'):
            source += f" - {chunk['source']}"
        if chunk.get('page_number'):
            source += f", p.{chunk['page_number']}"
        
        # Apply spacing fix to content
        fixed_content = fix_pdf_spacing(chunk['content'])
        context_parts.append(f"[Source: {source}]\n{fixed_content}")
    
    context = "\n\n---\n\n".join(context_parts)
    
    # Build prompt
    system_prompt = """You are a cardiology expert assistant. Answer questions based on the provided guideline excerpts.

FORMAT YOUR RESPONSE FOR READABILITY:
- Use **bold** for key terms, drug names, and important values
- Use bullet points (•) for lists of items
- Use clear section headers when covering multiple topics
- Highlight recommendation classes like **Class I**, **Class IIa**, **Class IIb**, **Class III**
- Highlight levels of evidence like **(LOE: A)**, **(LOE: B)**, **(LOE: C)**
- Use line breaks between sections for clarity

CONTENT RULES:
1. Base your answers on the provided guideline content
2. Cite specific guidelines (e.g., "Per the **2023 ACC/AHA AFib Guideline**...")
3. Include specific values: doses, percentages, timeframes, thresholds
4. If guidelines don't address the question, say so clearly
5. If there's conflicting information between guidelines, note this
6. Be thorough but organized"""

    user_prompt = f"""Based on the following guideline excerpts, please answer this question:

QUESTION: {question}

GUIDELINE EXCERPTS:
{context}

Please provide a clear, well-formatted, evidence-based answer with citations to the specific guidelines."""

    # Check if using OpenAI
    if model.startswith('gpt-'):
        return ask_openai(user_prompt, system_prompt, model, context_chunks)
    
    # Check if using extended thinking
    if model.endswith('-thinking'):
        actual_model = model.replace('-thinking', '')
        return ask_claude_thinking(user_prompt, system_prompt, actual_model, context_chunks)
    
    # Standard Claude call
    if not claude:
        return {'error': 'Anthropic API key not configured'}
    
    response = claude.messages.create(
        model=model,
        max_tokens=2000,
        system=system_prompt,
        messages=[{"role": "user", "content": user_prompt}]
    )
    
    # Calculate cost
    input_tokens = response.usage.input_tokens
    output_tokens = response.usage.output_tokens
    cost = calculate_cost(model, input_tokens, output_tokens)
    
    # Add summary to each source about what it covers
    def get_section_summary(content):
        """Extract a brief description of what this section covers."""
        content_lower = content.lower()
        topics = []
        
        # Check for common topics
        if any(word in content_lower for word in ['diagnosis', 'diagnostic', 'evaluation']):
            topics.append('Diagnosis')
        if any(word in content_lower for word in ['treatment', 'therapy', 'medication', 'drug']):
            topics.append('Treatment')
        if any(word in content_lower for word in ['class i', 'class ii', 'recommendation', 'recommended']):
            topics.append('Recommendations')
        if any(word in content_lower for word in ['epidemiology', 'prevalence', 'incidence', 'risk factor']):
            topics.append('Epidemiology')
        if any(word in content_lower for word in ['definition', 'criteria', 'classification']):
            topics.append('Definitions')
        if any(word in content_lower for word in ['surgery', 'surgical', 'intervention', 'procedure']):
            topics.append('Procedures')
        if any(word in content_lower for word in ['follow-up', 'monitoring', 'surveillance']):
            topics.append('Follow-up')
        if any(word in content_lower for word in ['prognosis', 'outcome', 'mortality']):
            topics.append('Prognosis')
        
        if topics:
            return ', '.join(topics[:3])
        return 'Clinical content'
    
    return {
        'answer': response.content[0].text,
        'model': model,
        'usage': {
            'input_tokens': input_tokens,
            'output_tokens': output_tokens,
            'total_tokens': input_tokens + output_tokens,
            'cost_usd': cost
        },
        'sources': [
            {
                'guideline': c['guideline_name'],
                'guideline_id': c['guideline_id'],
                'year': c.get('year'),
                'source': c.get('source'),
                'page': c.get('page_number'),
                'covers': get_section_summary(c['content']),
                'excerpt': fix_pdf_spacing(c['content'][:200]) + '...' if len(c['content']) > 200 else fix_pdf_spacing(c['content']),
                'full_text': fix_pdf_spacing(c['content'])
            }
            for c in context_chunks
        ]
    }

def ask_claude_thinking(user_prompt, system_prompt, model, context_chunks):
    """Ask Claude with extended thinking enabled."""
    if not claude:
        return {'error': 'Anthropic API key not configured'}
    
    try:
        # Extended thinking requires combining system + user prompt
        # because system parameter is not supported with thinking
        combined_prompt = f"""{system_prompt}

{user_prompt}"""
        
        response = claude.messages.create(
            model=model,
            max_tokens=16000,
            thinking={
                "type": "enabled",
                "budget_tokens": 10000
            },
            messages=[{"role": "user", "content": combined_prompt}]
        )
        
        # Extract the text response (skip thinking blocks)
        answer_text = ""
        for block in response.content:
            if block.type == "text":
                answer_text = block.text
                break
        
        if not answer_text:
            answer_text = "No response generated"
        
        # Calculate cost
        input_tokens = response.usage.input_tokens
        output_tokens = response.usage.output_tokens
        cost = calculate_cost(model, input_tokens, output_tokens)
        
        return {
            'answer': answer_text,
            'model': f"{model} (thinking)",
            'usage': {
                'input_tokens': input_tokens,
                'output_tokens': output_tokens,
                'total_tokens': input_tokens + output_tokens,
                'cost_usd': cost
            },
            'sources': [
                {
                    'guideline': c['guideline_name'],
                    'guideline_id': c['guideline_id'],
                    'year': c.get('year'),
                    'source': c.get('source'),
                    'page': c.get('page_number'),
                    'excerpt': fix_pdf_spacing(c['content'][:200]) + '...' if len(c['content']) > 200 else fix_pdf_spacing(c['content']),
                    'full_text': fix_pdf_spacing(c['content'])
                }
                for c in context_chunks
            ]
        }
    except Exception as e:
        import traceback
        print(f"Thinking mode error: {e}")
        print(f"Traceback: {traceback.format_exc()}")
        return {'error': f'Extended thinking failed: {str(e)}'}

def ask_openai(user_prompt, system_prompt, model, context_chunks):
    """Ask OpenAI GPT a question."""
    import requests
    
    openai_key = os.environ.get('OPENAI_API_KEY')
    if not openai_key:
        return {'error': 'OpenAI API key not configured'}
    
    try:
        # GPT-5.2 uses max_completion_tokens instead of max_tokens
        request_body = {
            'model': model,
            'messages': [
                {'role': 'system', 'content': system_prompt},
                {'role': 'user', 'content': user_prompt}
            ],
            'temperature': 0.1
        }
        
        # Use correct token parameter based on model
        if model.startswith('gpt-5'):
            request_body['max_completion_tokens'] = 2000
        else:
            request_body['max_tokens'] = 2000
        
        response = requests.post(
            'https://api.openai.com/v1/chat/completions',
            headers={
                'Authorization': f'Bearer {openai_key}',
                'Content-Type': 'application/json'
            },
            json=request_body,
            timeout=120
        )
        
        if response.status_code == 200:
            data = response.json()
            usage = data.get('usage', {})
            input_tokens = usage.get('prompt_tokens', 0)
            output_tokens = usage.get('completion_tokens', 0)
            cost = calculate_cost(model, input_tokens, output_tokens)
            
            # Reuse the section summary function 
            def get_section_summary(content):
                content_lower = content.lower()
                topics = []
                if any(word in content_lower for word in ['diagnosis', 'diagnostic', 'evaluation']):
                    topics.append('Diagnosis')
                if any(word in content_lower for word in ['treatment', 'therapy', 'medication', 'drug']):
                    topics.append('Treatment')
                if any(word in content_lower for word in ['class i', 'class ii', 'recommendation', 'recommended']):
                    topics.append('Recommendations')
                if any(word in content_lower for word in ['epidemiology', 'prevalence', 'incidence', 'risk factor']):
                    topics.append('Epidemiology')
                if any(word in content_lower for word in ['definition', 'criteria', 'classification']):
                    topics.append('Definitions')
                if any(word in content_lower for word in ['surgery', 'surgical', 'intervention', 'procedure']):
                    topics.append('Procedures')
                if any(word in content_lower for word in ['follow-up', 'monitoring', 'surveillance']):
                    topics.append('Follow-up')
                if any(word in content_lower for word in ['prognosis', 'outcome', 'mortality']):
                    topics.append('Prognosis')
                return ', '.join(topics[:3]) if topics else 'Clinical content'
            
            return {
                'answer': data['choices'][0]['message']['content'],
                'model': model,
                'usage': {
                    'input_tokens': input_tokens,
                    'output_tokens': output_tokens,
                    'total_tokens': input_tokens + output_tokens,
                    'cost_usd': cost
                },
                'sources': [
                    {
                        'guideline': c['guideline_name'],
                        'guideline_id': c['guideline_id'],
                        'year': c.get('year'),
                        'source': c.get('source'),
                        'page': c.get('page_number'),
                        'covers': get_section_summary(c['content']),
                        'excerpt': fix_pdf_spacing(c['content'][:200]) + '...' if len(c['content']) > 200 else fix_pdf_spacing(c['content']),
                        'full_text': fix_pdf_spacing(c['content'])
                    }
                    for c in context_chunks
                ]
            }
        else:
            return {'error': f'OpenAI API error: {response.status_code} - {response.text}'}
    except Exception as e:
        return {'error': f'OpenAI API error: {str(e)}'}

# =============================================================================
# API ROUTES
# =============================================================================

@app.route('/')
def index():
    """Serve the frontend."""
    return send_from_directory('.', 'index.html')

@app.route('/api/health', methods=['GET'])
def health():
    """Health check endpoint."""
    return jsonify({
        'status': 'healthy',
        'timestamp': datetime.now().isoformat(),
        'database': bool(DATABASE_URL),
        'anthropic': bool(ANTHROPIC_API_KEY),
        'openai': bool(os.environ.get('OPENAI_API_KEY'))
    })

@app.route('/api/guidelines', methods=['GET'])
def list_guidelines():
    """List all ingested guidelines with full metadata."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    cur.execute("""
        SELECT g.id, g.name, g.year, g.source, g.filename, g.file_hash,
               g.page_count, g.extracted_title, g.extracted_headers,
               g.file_size_bytes, g.created_at, g.updated_at,
               COUNT(c.id) as chunk_count
        FROM guidelines g
        LEFT JOIN chunks c ON g.id = c.guideline_id
        GROUP BY g.id
        ORDER BY g.created_at DESC
    """)
    
    guidelines = cur.fetchall()
    cur.close()
    conn.close()
    
    # Format for JSON
    for g in guidelines:
        if g['created_at']:
            g['created_at'] = g['created_at'].isoformat()
        if g['updated_at']:
            g['updated_at'] = g['updated_at'].isoformat()
        if g['file_size_bytes']:
            g['file_size_mb'] = round(g['file_size_bytes'] / (1024 * 1024), 2)
    
    return jsonify({'guidelines': guidelines})

@app.route('/api/guidelines/<int:guideline_id>', methods=['GET'])
def get_guideline(guideline_id):
    """Get detailed info about a specific guideline."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    cur.execute("""
        SELECT g.*, COUNT(c.id) as chunk_count
        FROM guidelines g
        LEFT JOIN chunks c ON g.id = c.guideline_id
        WHERE g.id = %s
        GROUP BY g.id
    """, (guideline_id,))
    
    guideline = cur.fetchone()
    cur.close()
    conn.close()
    
    if not guideline:
        return jsonify({'error': 'Guideline not found'}), 404
    
    if guideline['created_at']:
        guideline['created_at'] = guideline['created_at'].isoformat()
    if guideline['updated_at']:
        guideline['updated_at'] = guideline['updated_at'].isoformat()
    
    return jsonify(guideline)

@app.route('/api/upload', methods=['POST'])
def upload_pdfs():
    """Upload and ingest one or more PDF guidelines."""
    if 'files' not in request.files and 'file' not in request.files:
        return jsonify({'error': 'No files provided'}), 400
    
    # Handle both single file and multiple files
    files = request.files.getlist('files') or [request.files.get('file')]
    files = [f for f in files if f and f.filename]
    
    if not files:
        return jsonify({'error': 'No valid files provided'}), 400
    
    results = []
    
    for file in files:
        if not file.filename.lower().endswith('.pdf'):
            results.append({
                'filename': file.filename,
                'status': 'error',
                'message': 'File must be a PDF'
            })
            continue
        
        # Get metadata (use form data or defaults)
        # For multiple files, we auto-detect names
        name = request.form.get('name', '') if len(files) == 1 else ''
        if not name:
            name = file.filename.replace('.pdf', '').replace('_', ' ').replace('-', ' ')
        
        year = request.form.get('year')
        source = request.form.get('source', '')
        
        # Try to extract year from filename
        if not year:
            year_match = re.search(r'20\d{2}', file.filename)
            if year_match:
                year = int(year_match.group())
        
        if year:
            try:
                year = int(year)
            except ValueError:
                year = None
        
        # Save file temporarily
        filename = f"{hashlib.md5(file.filename.encode()).hexdigest()[:8]}_{file.filename}"
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        
        try:
            file.save(filepath)
            result = ingest_pdf(filepath, name, year, source)
            result['filename'] = file.filename
            results.append(result)
        except Exception as e:
            results.append({
                'filename': file.filename,
                'status': 'error',
                'message': str(e)
            })
        finally:
            # Clean up temp file
            if os.path.exists(filepath):
                os.remove(filepath)
    
    # Summary
    successful = sum(1 for r in results if r.get('status') == 'success')
    existing = sum(1 for r in results if r.get('status') == 'exists')
    failed = sum(1 for r in results if r.get('status') == 'error')
    
    return jsonify({
        'results': results,
        'summary': {
            'total': len(results),
            'successful': successful,
            'existing': existing,
            'failed': failed
        }
    })

@app.route('/api/ask', methods=['POST'])
def ask_question():
    """Ask a question about the guidelines."""
    data = request.get_json()
    
    if not data or 'question' not in data:
        return jsonify({'error': 'No question provided'}), 400
    
    question = data['question']
    model = data.get('model', 'claude-sonnet-4-20250514')
    
    # Search for relevant chunks - use SELECTED MODEL for query understanding
    chunks = search_guidelines(question, limit=50, model=model)
    
    # DEBUG: Log what we found
    print(f"=== ASK DEBUG ===")
    print(f"Question: {question}")
    print(f"Model: {model}")
    print(f"Chunks found: {len(chunks)}")
    for i, c in enumerate(chunks[:3]):
        print(f"Chunk {i+1} (page {c.get('page_number')}): {c['content'][:100]}...")
    print(f"=================")
    
    if not chunks:
        return jsonify({
            'answer': "I couldn't find any relevant information in the guidelines database. Please make sure guidelines have been uploaded and try rephrasing your question.",
            'sources': []
        })
    
    # Get answer from Claude/GPT
    result = ask_claude(question, chunks, model)
    
    return jsonify(result)

@app.route('/api/delete/<int:guideline_id>', methods=['DELETE'])
def delete_guideline(guideline_id):
    """Delete a guideline and all its chunks."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Get guideline info first
    cur.execute("SELECT name, filename FROM guidelines WHERE id = %s", (guideline_id,))
    guideline = cur.fetchone()
    
    if not guideline:
        cur.close()
        conn.close()
        return jsonify({'error': 'Guideline not found'}), 404
    
    # Count chunks that will be deleted
    cur.execute("SELECT COUNT(*) as count FROM chunks WHERE guideline_id = %s", (guideline_id,))
    chunk_count = cur.fetchone()['count']
    
    # Delete (chunks will cascade)
    cur.execute("DELETE FROM guidelines WHERE id = %s", (guideline_id,))
    
    conn.commit()
    cur.close()
    conn.close()
    
    return jsonify({
        'status': 'deleted',
        'name': guideline['name'],
        'filename': guideline['filename'],
        'chunks_deleted': chunk_count
    })

@app.route('/api/stats', methods=['GET'])
def get_stats():
    """Get database statistics."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    cur.execute("""
        SELECT 
            COUNT(DISTINCT g.id) as total_guidelines,
            COUNT(c.id) as total_chunks,
            SUM(g.page_count) as total_pages,
            SUM(g.file_size_bytes) as total_size_bytes
        FROM guidelines g
        LEFT JOIN chunks c ON g.id = c.guideline_id
    """)
    
    stats = cur.fetchone()
    cur.close()
    conn.close()
    
    if stats['total_size_bytes']:
        stats['total_size_mb'] = round(stats['total_size_bytes'] / (1024 * 1024), 2)
    
    return jsonify(stats)

# =============================================================================
# TOPIC EXTRACTION & MANAGEMENT
# =============================================================================

@app.route('/api/topics/<int:guideline_id>')
def get_guideline_topics(guideline_id):
    """Get saved topics for a guideline, or return empty if not analyzed yet."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Check if topics already exist in database
    cur.execute("""
        SELECT id, topic_name, icon, start_page, end_page, chunk_ids
        FROM guideline_topics 
        WHERE guideline_id = %s
        ORDER BY start_page, id
    """, (guideline_id,))
    saved_topics = cur.fetchall()
    
    # Get guideline name
    cur.execute("SELECT name FROM guidelines WHERE id = %s", (guideline_id,))
    guideline = cur.fetchone()
    
    cur.close()
    conn.close()
    
    if not guideline:
        return jsonify({'error': 'Guideline not found'}), 404
    
    if saved_topics:
        # Return saved topics
        topics = [{'topic': t['topic_name'], 'icon': t['icon'], 'id': t['id'], 
                   'start_page': t['start_page'], 'end_page': t['end_page']} 
                  for t in saved_topics]
        return jsonify({
            'topics': topics, 
            'guideline': guideline['name'],
            'analyzed': True
        })
    else:
        # No topics saved yet
        return jsonify({
            'topics': [], 
            'guideline': guideline['name'],
            'analyzed': False,
            'message': 'Click "Analyze Topics" to extract topics from this guideline'
        })

@app.route('/api/analyze-topics/<int:guideline_id>', methods=['POST'])
def analyze_guideline_topics(guideline_id):
    """Use AI to analyze and save topics for a guideline."""
    import json
    
    # Get model from request body (default to Sonnet)
    data = request.get_json() or {}
    model = data.get('model', 'claude-sonnet-4-20250514')
    
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Get guideline info
    cur.execute("SELECT id, name, extracted_headers, page_count FROM guidelines WHERE id = %s", (guideline_id,))
    guideline = cur.fetchone()
    
    if not guideline:
        cur.close()
        conn.close()
        return jsonify({'error': 'Guideline not found'}), 404
    
    # Delete any existing topics for this guideline
    cur.execute("DELETE FROM guideline_topics WHERE guideline_id = %s", (guideline_id,))
    
    # Get ALL chunks with page numbers (excluding references)
    cur.execute("""
        SELECT id, content, page_number 
        FROM chunks 
        WHERE guideline_id = %s 
        AND (is_reference IS NULL OR is_reference = FALSE)
        ORDER BY page_number
    """, (guideline_id,))
    all_chunks = cur.fetchall()
    
    if not all_chunks:
        cur.close()
        conn.close()
        return jsonify({'error': 'No content found for this guideline'}), 404
    
    print(f"=== TOPIC ANALYSIS ===")
    print(f"Guideline: {guideline['name']}")
    print(f"Total chunks: {len(all_chunks)}")
    
    # Build content for AI analysis - sample from throughout the document
    headers = guideline.get('extracted_headers', '') or ''
    
    # Get chunks spread throughout document
    step = max(1, len(all_chunks) // 20)  # Sample ~20 chunks
    sample_chunks = all_chunks[::step][:25]
    content_sample = '\n\n'.join([
        f"[Page {c['page_number']}]: {c['content'][:600]}" 
        for c in sample_chunks
    ])
    
    if not claude:
        cur.close()
        conn.close()
        return jsonify({'error': 'Anthropic API key not configured'}), 500
    
    # AI prompt to extract topics with page ranges
    prompt = f"""Analyze this medical guideline and identify ALL major clinical sections/topics.

GUIDELINE: {guideline['name']}
TOTAL PAGES: {guideline['page_count']}

EXTRACTED HEADERS FROM PDF:
{headers[:3000]}

SAMPLE CONTENT FROM THROUGHOUT DOCUMENT:
{content_sample[:6000]}

Identify 5-15 distinct clinical TOPICS/SECTIONS that would each make a good standalone presentation.

For each topic, estimate the page range where it appears.

DO NOT include these as topics:
- References / Bibliography
- Authors / Contributors  
- Disclosures / Conflicts of Interest
- Acknowledgments
- Preamble / Introduction (unless it has clinical content)
- Methods

Return JSON array:
[
    {{
        "topic": "Clear topic name (e.g., 'Anticoagulation Therapy', 'Surgical Indications')",
        "icon": "relevant emoji",
        "start_page": estimated_start_page,
        "end_page": estimated_end_page,
        "description": "Brief 1-line description of what this section covers"
    }}
]

Use medical emojis: 💊 💉 🫀 ❤️ 🧠 ⚡ 🔥 🩺 📊 🎯 ⚠️ 🛡️ 💓 🩸 🔬 📋 💪 🏥

Return ONLY valid JSON array."""

    try:
        print(f"Using model for topic analysis: {model}")
        
        if model.startswith('gpt-'):
            # Use OpenAI for GPT models
            import requests as req
            openai_key = os.environ.get('OPENAI_API_KEY')
            resp = req.post(
                'https://api.openai.com/v1/chat/completions',
                headers={'Authorization': f'Bearer {openai_key}', 'Content-Type': 'application/json'},
                json={
                    'model': model,
                    'messages': [{'role': 'user', 'content': prompt}],
                    'temperature': 0,
                    'max_completion_tokens': 2000
                },
                timeout=120
            )
            if resp.status_code == 200:
                result_text = resp.json()['choices'][0]['message']['content'].strip()
            else:
                print(f"OpenAI topic analysis error: {resp.status_code} - {resp.text[:500]}")
                return jsonify({'error': f'OpenAI API error: {resp.status_code}'}), 500
        else:
            # Use Claude
            response = claude.messages.create(
                model=model,
                max_tokens=2000,
                temperature=0,
                messages=[{"role": "user", "content": prompt}]
            )
            result_text = response.content[0].text.strip()
        
        # Parse JSON
        start = result_text.find('[')
        end = result_text.rfind(']') + 1
        if start >= 0 and end > start:
            topics = json.loads(result_text[start:end])
            
            # Save topics to database
            saved_topics = []
            for t in topics:
                topic_name = t.get('topic', 'Unknown')
                icon = t.get('icon', '📋')
                start_page = t.get('start_page')
                end_page = t.get('end_page')
                
                # Find chunk IDs in this page range
                chunk_ids = [str(c['id']) for c in all_chunks 
                            if start_page and end_page and 
                            start_page <= c['page_number'] <= end_page]
                
                cur.execute("""
                    INSERT INTO guideline_topics (guideline_id, topic_name, icon, start_page, end_page, chunk_ids)
                    VALUES (%s, %s, %s, %s, %s, %s)
                    RETURNING id
                """, (guideline_id, topic_name, icon, start_page, end_page, ','.join(chunk_ids) if chunk_ids else None))
                
                topic_id = cur.fetchone()['id']
                saved_topics.append({
                    'id': topic_id,
                    'topic': topic_name,
                    'icon': icon,
                    'start_page': start_page,
                    'end_page': end_page,
                    'chunk_count': len(chunk_ids)
                })
            
            conn.commit()
            cur.close()
            conn.close()
            
            print(f"Saved {len(saved_topics)} topics for guideline {guideline_id}")
            
            return jsonify({
                'success': True,
                'topics': saved_topics,
                'guideline': guideline['name']
            })
        else:
            cur.close()
            conn.close()
            return jsonify({'error': 'Failed to parse AI response'}), 500
            
    except Exception as e:
        print(f"Topic analysis error: {e}")
        import traceback
        print(traceback.format_exc())
        conn.rollback()
        cur.close()
        conn.close()
        return jsonify({'error': str(e)}), 500

def extract_topics_from_headers(headers):
    """Fallback: extract topics from headers if AI fails."""
    if not headers:
        return []
    
    # Common clinical topic patterns
    topic_keywords = [
        'anticoagulation', 'ablation', 'rate control', 'rhythm control',
        'stroke', 'heart failure', 'prevention', 'diagnosis', 'treatment',
        'management', 'therapy', 'screening', 'risk', 'assessment'
    ]
    
    topics = []
    seen = set()
    
    for line in headers.split('\n'):
        line_lower = line.lower().strip()
        if len(line_lower) < 5 or len(line_lower) > 60:
            continue
        if any(skip in line_lower for skip in ['reference', 'author', 'disclosure', 'appendix', 'acknowledge']):
            continue
        if any(kw in line_lower for kw in topic_keywords):
            topic_name = line.strip()
            if topic_name.lower() not in seen:
                seen.add(topic_name.lower())
                topics.append({'topic': topic_name, 'icon': '📋'})
    
    return topics[:10]

# =============================================================================
# PDF PAGE IMAGE EXTRACTION
# =============================================================================

@app.route('/api/page-image/<int:guideline_id>/<int:page_number>')
def get_page_image(guideline_id, page_number):
    """Extract and return a specific page from a stored PDF as an image."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Get PDF data
    cur.execute("SELECT pdf_data, page_count, name FROM guidelines WHERE id = %s", (guideline_id,))
    guideline = cur.fetchone()
    cur.close()
    conn.close()
    
    if not guideline:
        return jsonify({'error': 'Guideline not found'}), 404
    
    if not guideline['pdf_data']:
        return jsonify({'error': 'PDF not stored for this guideline'}), 404
    
    if page_number < 1 or page_number > guideline['page_count']:
        return jsonify({'error': f'Invalid page number. PDF has {guideline["page_count"]} pages'}), 400
    
    try:
        # Try to use pdf2image (requires poppler)
        try:
            from pdf2image import convert_from_bytes
            import io
            
            # Convert just the requested page (page_number is 1-indexed)
            images = convert_from_bytes(
                bytes(guideline['pdf_data']),
                first_page=page_number,
                last_page=page_number,
                dpi=150  # Good balance of quality and size
            )
            
            if images:
                # Convert to PNG
                img_buffer = io.BytesIO()
                images[0].save(img_buffer, format='PNG')
                img_buffer.seek(0)
                
                from flask import send_file
                return send_file(
                    img_buffer,
                    mimetype='image/png',
                    download_name=f'{guideline["name"]}_page_{page_number}.png'
                )
        except ImportError:
            # pdf2image not available, try pymupdf
            try:
                import fitz  # PyMuPDF
                import io
                
                pdf_doc = fitz.open(stream=bytes(guideline['pdf_data']), filetype="pdf")
                page = pdf_doc[page_number - 1]  # 0-indexed
                
                # Render at 2x for good quality
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                
                img_buffer = io.BytesIO(pix.tobytes("png"))
                img_buffer.seek(0)
                
                from flask import send_file
                return send_file(
                    img_buffer,
                    mimetype='image/png',
                    download_name=f'{guideline["name"]}_page_{page_number}.png'
                )
            except ImportError:
                return jsonify({'error': 'PDF image extraction not available. Install pdf2image or PyMuPDF'}), 500
                
    except Exception as e:
        print(f"Page image extraction error: {e}")
        import traceback
        print(traceback.format_exc())
        return jsonify({'error': str(e)}), 500

# =============================================================================
# PPTX GENERATION - BACKGROUND JOB SYSTEM
# =============================================================================

import threading

@app.route('/api/generate-pptx/<int:guideline_id>', methods=['POST'])
def start_pptx_job(guideline_id):
    """Start a background job to generate PPTX. Returns job ID for polling."""
    import json
    
    data = request.get_json() or {}
    topic = data.get('topic')
    model = data.get('model', 'claude-sonnet-4-20250514')
    
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Check guideline exists
    cur.execute("SELECT id, name FROM guidelines WHERE id = %s", (guideline_id,))
    guideline = cur.fetchone()
    
    if not guideline:
        cur.close()
        conn.close()
        return jsonify({'error': 'Guideline not found'}), 404
    
    # Count chunks to estimate batches - use page ranges if available
    if topic:
        # Check for saved topic page ranges
        cur.execute("""
            SELECT start_page, end_page 
            FROM guideline_topics 
            WHERE guideline_id = %s AND topic_name = %s
        """, (guideline_id, topic))
        topic_info = cur.fetchone()
        
        if topic_info and topic_info['start_page'] and topic_info['end_page']:
            cur.execute("""
                SELECT COUNT(*) as cnt FROM chunks 
                WHERE guideline_id = %s 
                AND page_number >= %s AND page_number <= %s
                AND (is_reference IS NULL OR is_reference = FALSE)
            """, (guideline_id, topic_info['start_page'], topic_info['end_page']))
        else:
            cur.execute("""
                SELECT COUNT(*) as cnt FROM chunks 
                WHERE guideline_id = %s 
                AND LOWER(content) LIKE %s
                AND (is_reference IS NULL OR is_reference = FALSE)
            """, (guideline_id, f'%{topic.lower()}%'))
    else:
        cur.execute("""
            SELECT COUNT(*) as cnt FROM chunks 
            WHERE guideline_id = %s
            AND (is_reference IS NULL OR is_reference = FALSE)
        """, (guideline_id,))
    
    chunk_count = cur.fetchone()['cnt']
    total_batches = (chunk_count + 24) // 25  # 25 chunks per batch
    
    # Create job record
    cur.execute("""
        INSERT INTO pptx_jobs (guideline_id, topic, model, status, total_batches)
        VALUES (%s, %s, %s, 'processing', %s)
        RETURNING id
    """, (guideline_id, topic, model, total_batches))
    job_id = cur.fetchone()['id']
    conn.commit()
    cur.close()
    conn.close()
    
    # Start background thread
    thread = threading.Thread(target=process_pptx_job, args=(job_id,))
    thread.daemon = True
    thread.start()
    
    return jsonify({
        'job_id': job_id,
        'status': 'processing',
        'total_batches': total_batches,
        'estimated_time': f"{total_batches * 5}-{total_batches * 8} seconds"
    })

@app.route('/api/pptx-job/<int:job_id>')
def get_pptx_job_status(job_id):
    """Poll job status."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    cur.execute("""
        SELECT j.*, g.name as guideline_name 
        FROM pptx_jobs j
        JOIN guidelines g ON j.guideline_id = g.id
        WHERE j.id = %s
    """, (job_id,))
    job = cur.fetchone()
    cur.close()
    conn.close()
    
    if not job:
        return jsonify({'error': 'Job not found'}), 404
    
    return jsonify({
        'job_id': job['id'],
        'status': job['status'],
        'progress': job['progress'],
        'total_batches': job['total_batches'],
        'cost': job['cost'],
        'tokens': job['tokens'],
        'error': job['error'],
        'guideline_name': job['guideline_name'],
        'topic': job['topic']
    })

@app.route('/api/library')
def get_library():
    """Get all completed jobs grouped by guideline for the Library tab."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    cur.execute("""
        SELECT j.id, j.guideline_id, j.topic, j.model, j.cost, j.tokens, 
               j.created_at, j.completed_at,
               j.pptx_data IS NOT NULL as has_pptx_cache,
               j.html_data IS NOT NULL as has_html_cache,
               g.name as guideline_name, g.year, g.source
        FROM pptx_jobs j
        JOIN guidelines g ON j.guideline_id = g.id
        WHERE j.status = 'completed'
        ORDER BY g.name, j.completed_at DESC
    """)
    jobs = cur.fetchall()
    cur.close()
    conn.close()
    
    # Group by guideline
    guidelines = {}
    for job in jobs:
        g_id = job['guideline_id']
        if g_id not in guidelines:
            guidelines[g_id] = {
                'id': g_id,
                'name': job['guideline_name'],
                'year': job['year'],
                'source': job['source'],
                'jobs': []
            }
        guidelines[g_id]['jobs'].append({
            'id': job['id'],
            'topic': job['topic'] or 'Full Guideline',
            'model': job['model'],
            'cost': job['cost'],
            'tokens': job['tokens'],
            'created_at': job['created_at'].isoformat() if job['created_at'] else None,
            'completed_at': job['completed_at'].isoformat() if job['completed_at'] else None,
            'has_pptx_cache': job['has_pptx_cache'],
            'has_html_cache': job['has_html_cache']
        })
    
    return jsonify({
        'guidelines': list(guidelines.values()),
        'total_jobs': len(jobs)
    })

@app.route('/api/check-spacing-issues')
def check_spacing_issues():
    """Check all chunks for spacing issues and report affected guidelines."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Get all chunks with their guideline info
    cur.execute("""
        SELECT c.id, c.content, c.guideline_id, g.name as guideline_name
        FROM chunks c
        JOIN guidelines g ON c.guideline_id = g.id
    """)
    
    chunks = cur.fetchall()
    cur.close()
    conn.close()
    
    # Check each chunk for spacing issues
    issues_by_guideline = {}
    total_bad_chunks = 0
    
    for chunk in chunks:
        content = chunk['content']
        # Check for long words (likely missing spaces)
        words = content.split()
        long_words = [w for w in words if len(w) > 35 and re.search(r'[a-z]{15,}', w)]
        
        if long_words:
            total_bad_chunks += 1
            gid = chunk['guideline_id']
            gname = chunk['guideline_name']
            
            if gid not in issues_by_guideline:
                issues_by_guideline[gid] = {
                    'id': gid,
                    'name': gname,
                    'bad_chunks': 0,
                    'examples': []
                }
            
            issues_by_guideline[gid]['bad_chunks'] += 1
            
            # Keep up to 3 examples per guideline
            if len(issues_by_guideline[gid]['examples']) < 3:
                issues_by_guideline[gid]['examples'].append({
                    'chunk_id': chunk['id'],
                    'sample': long_words[0][:80] if long_words else ''
                })
    
    return jsonify({
        'total_chunks': len(chunks),
        'bad_chunks': total_bad_chunks,
        'affected_guidelines': list(issues_by_guideline.values())
    })

@app.route('/api/fix-guideline-spacing/<int:guideline_id>', methods=['POST'])
def fix_guideline_spacing(guideline_id):
    """Start a background job to fix spacing issues in a guideline's chunks."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Get guideline
    cur.execute("SELECT id, name FROM guidelines WHERE id = %s", (guideline_id,))
    guideline = cur.fetchone()
    
    if not guideline:
        cur.close()
        conn.close()
        return jsonify({'error': 'Guideline not found'}), 404
    
    # Count chunks that need fixing
    cur.execute("SELECT id, content FROM chunks WHERE guideline_id = %s", (guideline_id,))
    all_chunks = cur.fetchall()
    
    chunks_to_fix = [c for c in all_chunks if check_text_needs_spacing_fix(c['content'])]
    
    if not chunks_to_fix:
        cur.close()
        conn.close()
        return jsonify({
            'success': True,
            'message': 'No chunks need fixing',
            'total_chunks': len(all_chunks),
            'bad_chunks': 0
        })
    
    # Create a background job
    cur.execute("""
        INSERT INTO spacing_jobs (guideline_id, status, total_chunks, progress)
        VALUES (%s, 'pending', %s, 0)
        RETURNING id
    """, (guideline_id, len(chunks_to_fix)))
    job_id = cur.fetchone()['id']
    conn.commit()
    cur.close()
    conn.close()
    
    # Start background processing
    import threading
    thread = threading.Thread(target=process_spacing_job, args=(job_id, guideline_id, chunks_to_fix))
    thread.daemon = True
    thread.start()
    
    return jsonify({
        'success': True,
        'job_id': job_id,
        'guideline': guideline['name'],
        'total_chunks': len(all_chunks),
        'chunks_to_fix': len(chunks_to_fix),
        'message': f'Started fixing {len(chunks_to_fix)} chunks in background'
    })

def process_spacing_job(job_id, guideline_id, chunks_to_fix):
    """Background worker to fix spacing in chunks using batched AI calls."""
    print(f"[Spacing Job {job_id}] Starting - {len(chunks_to_fix)} chunks to fix")
    
    BATCH_SIZE = 5  # Process 5 chunks per AI call
    
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Update status to processing
    cur.execute("UPDATE spacing_jobs SET status = 'processing' WHERE id = %s", (job_id,))
    conn.commit()
    
    fixed_count = 0
    processed_count = 0
    
    # Process in batches
    for batch_start in range(0, len(chunks_to_fix), BATCH_SIZE):
        batch = chunks_to_fix[batch_start:batch_start + BATCH_SIZE]
        
        try:
            # Build batch prompt
            batch_texts = []
            for i, chunk in enumerate(batch):
                batch_texts.append(f"[CHUNK {i+1}]\n{chunk['content']}\n[/CHUNK {i+1}]")
            
            combined_text = "\n\n".join(batch_texts)
            
            # Single AI call for entire batch
            if claude:
                response = claude.messages.create(
                    model="claude-sonnet-4-20250514",
                    max_tokens=len(combined_text) + 2000,
                    messages=[{
                        "role": "user",
                        "content": f"""Fix the spacing in these text chunks extracted from a PDF. Words are concatenated without spaces.

RULES:
1. Add spaces where words run together
2. Keep all original words and meaning
3. Preserve the [CHUNK N] markers exactly
4. Return ALL chunks with fixed spacing

{combined_text}

Return the fixed chunks with the same [CHUNK N] markers:"""
                    }]
                )
                
                result = response.content[0].text
                
                # Parse results and update each chunk
                for i, chunk in enumerate(batch):
                    chunk_num = i + 1
                    # Extract this chunk's fixed content
                    pattern = rf'\[CHUNK {chunk_num}\]\s*(.*?)\s*\[/CHUNK {chunk_num}\]'
                    match = re.search(pattern, result, re.DOTALL)
                    
                    if match:
                        fixed_content = match.group(1).strip()
                        if fixed_content and fixed_content != chunk['content']:
                            cur.execute(
                                "UPDATE chunks SET content = %s WHERE id = %s",
                                (fixed_content, chunk['id'])
                            )
                            fixed_count += 1
                    
                    processed_count += 1
                
                print(f"[Spacing Job {job_id}] Batch {batch_start//BATCH_SIZE + 1}: fixed {fixed_count} so far ({processed_count}/{len(chunks_to_fix)})")
            else:
                # No AI - use DP fallback for each chunk
                for chunk in batch:
                    fixed_content = fix_text_spacing(chunk['content'])
                    if fixed_content != chunk['content']:
                        cur.execute(
                            "UPDATE chunks SET content = %s WHERE id = %s",
                            (fixed_content, chunk['id'])
                        )
                        fixed_count += 1
                    processed_count += 1
                    
        except Exception as e:
            print(f"[Spacing Job {job_id}] Batch failed: {e}, using DP fallback")
            # Fall back to DP algorithm for this batch
            for chunk in batch:
                try:
                    fixed_content = fix_text_spacing(chunk['content'])
                    if fixed_content != chunk['content']:
                        cur.execute(
                            "UPDATE chunks SET content = %s WHERE id = %s",
                            (fixed_content, chunk['id'])
                        )
                        fixed_count += 1
                except:
                    pass
                processed_count += 1
        
        # Update progress after each batch
        progress = int(processed_count / len(chunks_to_fix) * 100)
        cur.execute(
            "UPDATE spacing_jobs SET progress = %s, fixed_chunks = %s WHERE id = %s",
            (progress, fixed_count, job_id)
        )
        conn.commit()
    
    # Mark as completed
    cur.execute("""
        UPDATE spacing_jobs 
        SET status = 'completed', progress = 100, fixed_chunks = %s, completed_at = NOW()
        WHERE id = %s
    """, (fixed_count, job_id))
    conn.commit()
    
    cur.close()
    conn.close()
    
    print(f"[Spacing Job {job_id}] Completed - fixed {fixed_count} of {len(chunks_to_fix)} chunks")

@app.route('/api/spacing-job/<int:job_id>')
def get_spacing_job_status(job_id):
    """Get the status of a spacing fix job."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    cur.execute("""
        SELECT sj.*, g.name as guideline_name
        FROM spacing_jobs sj
        JOIN guidelines g ON sj.guideline_id = g.id
        WHERE sj.id = %s
    """, (job_id,))
    job = cur.fetchone()
    
    cur.close()
    conn.close()
    
    if not job:
        return jsonify({'error': 'Job not found'}), 404
    
    return jsonify({
        'id': job['id'],
        'guideline': job['guideline_name'],
        'status': job['status'],
        'progress': job['progress'],
        'total_chunks': job['total_chunks'],
        'fixed_chunks': job['fixed_chunks'],
        'error': job['error'],
        'created_at': job['created_at'].isoformat() if job['created_at'] else None,
        'completed_at': job['completed_at'].isoformat() if job['completed_at'] else None
    })

@app.route('/api/pptx-job/<int:job_id>', methods=['DELETE'])
def delete_pptx_job(job_id):
    """Delete a generated PPTX/HTML job."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    # Check if job exists
    cur.execute("SELECT id, guideline_id, topic FROM pptx_jobs WHERE id = %s", (job_id,))
    job = cur.fetchone()
    
    if not job:
        cur.close()
        conn.close()
        return jsonify({'error': 'Job not found'}), 404
    
    # Delete the job
    cur.execute("DELETE FROM pptx_jobs WHERE id = %s", (job_id,))
    conn.commit()
    
    cur.close()
    conn.close()
    
    print(f"Deleted job {job_id} (topic: {job['topic']})")
    
    return jsonify({
        'success': True,
        'message': f"Deleted job {job_id}",
        'topic': job['topic']
    })

@app.route('/api/pptx-job/<int:job_id>/download')
def download_pptx_job(job_id):
    """Download completed PPTX - uses cached version if available."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    cur.execute("""
        SELECT j.*, g.name as guideline_name, g.year, g.source
        FROM pptx_jobs j
        JOIN guidelines g ON j.guideline_id = g.id
        WHERE j.id = %s AND j.status = 'completed'
    """, (job_id,))
    job = cur.fetchone()
    
    if not job:
        cur.close()
        conn.close()
        return jsonify({'error': 'Job not found or not completed'}), 404
    
    if not job['result_data']:
        cur.close()
        conn.close()
        return jsonify({'error': 'No result data'}), 404
    
    filename = f"{job['guideline_name'].replace(' ', '_')}{'_' + job['topic'] if job['topic'] else ''}.pptx"
    
    # Check if we have cached PPTX data
    if job.get('pptx_data'):
        print(f"Serving cached PPTX for job {job_id}")
        cur.close()
        conn.close()
        from flask import send_file, make_response
        import io
        pptx_buffer = io.BytesIO(bytes(job['pptx_data']))
        response = make_response(send_file(
            pptx_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        ))
        response.headers['X-Generation-Cost'] = f"{job['cost']:.4f}" if job['cost'] else '0'
        response.headers['X-Generation-Tokens'] = str(job['tokens'] or 0)
        response.headers['X-Cached'] = 'true'
        response.headers['Access-Control-Expose-Headers'] = 'X-Generation-Cost, X-Generation-Tokens, X-Cached'
        return response
    
    # Generate PPTX and cache it
    print(f"Generating new PPTX for job {job_id}")
    import json
    summary_json = json.loads(job['result_data'])
    
    guideline = {
        'id': job['guideline_id'],
        'name': job['guideline_name'],
        'year': job['year'],
        'source': job['source']
    }
    
    pptx_path = create_pptx(summary_json, guideline, job['topic'])
    
    if pptx_path:
        # Read and cache the PPTX file
        with open(pptx_path, 'rb') as f:
            pptx_bytes = f.read()
        
        # Save to database for future downloads
        cur.execute("UPDATE pptx_jobs SET pptx_data = %s WHERE id = %s", (pptx_bytes, job_id))
        conn.commit()
        print(f"Cached PPTX for job {job_id} ({len(pptx_bytes)} bytes)")
        
        cur.close()
        conn.close()
        
        from flask import send_file, make_response
        response = make_response(send_file(
            pptx_path,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        ))
        response.headers['X-Generation-Cost'] = f"{job['cost']:.4f}" if job['cost'] else '0'
        response.headers['X-Generation-Tokens'] = str(job['tokens'] or 0)
        response.headers['X-Cached'] = 'false'
        response.headers['Access-Control-Expose-Headers'] = 'X-Generation-Cost, X-Generation-Tokens, X-Cached'
        return response
    else:
        cur.close()
        conn.close()
        return jsonify({'error': 'Failed to create PPTX'}), 500

@app.route('/api/pptx-job/<int:job_id>/download-html')
def download_pptx_job_html(job_id):
    """Download completed extraction as beautiful HTML - uses cached version if available."""
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    cur.execute("""
        SELECT j.*, g.name as guideline_name, g.year, g.source
        FROM pptx_jobs j
        JOIN guidelines g ON j.guideline_id = g.id
        WHERE j.id = %s AND j.status = 'completed'
    """, (job_id,))
    job = cur.fetchone()
    
    if not job:
        cur.close()
        conn.close()
        return jsonify({'error': 'Job not found or not completed'}), 404
    
    if not job['result_data']:
        cur.close()
        conn.close()
        return jsonify({'error': 'No result data'}), 404
    
    filename = f"{job['guideline_name'].replace(' ', '_')}{'_' + job['topic'] if job['topic'] else ''}.html"
    
    # Check if we have cached HTML data
    if job.get('html_data'):
        print(f"Serving cached HTML for job {job_id}")
        cur.close()
        conn.close()
        from flask import Response
        return Response(
            job['html_data'],
            mimetype='text/html',
            headers={'Content-Disposition': f'attachment; filename="{filename}"'}
        )
    
    # Generate HTML
    print(f"Generating new HTML for job {job_id}")
    import json
    from flask import Response
    
    summary = json.loads(job['result_data'])
    
    section_config = {
        'recommendations': {'title': 'Clinical Recommendations', 'icon': '⚕️', 'color': '#22c55e'},
        'definitions': {'title': 'Definitions & Classifications', 'icon': '📖', 'color': '#3b82f6'},
        'epidemiology': {'title': 'Epidemiology & Risk Factors', 'icon': '📊', 'color': '#a855f7'},
        'pathophysiology': {'title': 'Pathophysiology', 'icon': '🔬', 'color': '#ec4899'},
        'diagnosis': {'title': 'Diagnosis & Evaluation', 'icon': '🩺', 'color': '#f59e0b'},
        'treatment': {'title': 'Treatment & Therapy', 'icon': '💊', 'color': '#14b8a6'},
        'management': {'title': 'Management & Follow-up', 'icon': '📋', 'color': '#f97316'},
        'special_populations': {'title': 'Special Populations', 'icon': '👥', 'color': '#8b5cf6'},
        'key_points': {'title': 'Key Clinical Points', 'icon': '⭐', 'color': '#ef4444'},
    }
    
    # Build HTML content
    html = f'''<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{job['guideline_name']}{' - ' + job['topic'] if job['topic'] else ''}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%);
            color: #e2e8f0;
            line-height: 1.7;
            min-height: 100vh;
            padding: 40px 20px;
        }}
        .container {{ max-width: 900px; margin: 0 auto; }}
        
        .header {{
            text-align: center;
            padding: 40px;
            background: linear-gradient(135deg, #1e293b 0%, #334155 100%);
            border-radius: 20px;
            margin-bottom: 30px;
            border: 1px solid #475569;
        }}
        .header h1 {{ font-size: 2rem; margin-bottom: 10px; color: #fff; }}
        .header .topic {{ font-size: 1.3rem; color: #60a5fa; margin-bottom: 15px; }}
        .header .meta {{ font-size: 0.9rem; color: #94a3b8; }}
        
        .toc {{
            background: #1e293b;
            padding: 25px;
            border-radius: 15px;
            margin-bottom: 30px;
            border: 1px solid #334155;
        }}
        .toc h2 {{ font-size: 1.1rem; margin-bottom: 15px; color: #94a3b8; }}
        .toc-item {{
            display: inline-block;
            padding: 8px 16px;
            margin: 5px;
            background: #334155;
            border-radius: 8px;
            color: #e2e8f0;
            text-decoration: none;
            font-size: 0.9rem;
            transition: all 0.2s;
        }}
        .toc-item:hover {{ background: #475569; transform: translateY(-2px); }}
        
        .section {{
            background: #1e293b;
            border-radius: 15px;
            margin-bottom: 25px;
            overflow: hidden;
            border: 1px solid #334155;
        }}
        .section-header {{
            padding: 20px 25px;
            display: flex;
            align-items: center;
            gap: 12px;
            cursor: pointer;
            user-select: none;
        }}
        .section-header:hover {{ background: rgba(255,255,255,0.05); }}
        .section-icon {{ font-size: 1.5rem; }}
        .section-title {{ font-size: 1.2rem; font-weight: 600; flex: 1; }}
        .section-count {{
            background: rgba(255,255,255,0.1);
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 0.85rem;
        }}
        .section-toggle {{ font-size: 1.2rem; transition: transform 0.3s; }}
        .section.collapsed .section-toggle {{ transform: rotate(-90deg); }}
        .section.collapsed .section-content {{ display: none; }}
        
        .section-content {{ padding: 0 25px 25px 25px; }}
        
        .item {{
            background: #0f172a;
            padding: 20px;
            border-radius: 10px;
            margin-bottom: 15px;
            border-left: 4px solid;
        }}
        .item-text {{ margin-bottom: 10px; }}
        .item-meta {{
            display: flex;
            gap: 10px;
            flex-wrap: wrap;
        }}
        .badge {{
            padding: 4px 10px;
            border-radius: 5px;
            font-size: 0.8rem;
            font-weight: 500;
        }}
        .badge-page {{ background: #334155; color: #94a3b8; }}
        .badge-class-i {{ background: rgba(34,197,94,0.2); color: #22c55e; }}
        .badge-class-iia {{ background: rgba(59,130,246,0.2); color: #3b82f6; }}
        .badge-class-iib {{ background: rgba(251,191,36,0.2); color: #fbbf24; }}
        .badge-class-iii {{ background: rgba(239,68,68,0.2); color: #ef4444; }}
        .badge-loe {{ background: rgba(168,85,247,0.2); color: #a855f7; }}
        
        .footer {{
            text-align: center;
            padding: 30px;
            color: #64748b;
            font-size: 0.85rem;
        }}
        
        @media print {{
            body {{ background: #fff; color: #000; padding: 20px; }}
            .section {{ break-inside: avoid; border: 1px solid #ccc; }}
            .section-header {{ background: #f0f0f0 !important; }}
            .item {{ background: #f8f8f8; border-left-width: 3px; }}
            .toc {{ background: #f0f0f0; }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <header class="header">
            <h1>{job['guideline_name']}</h1>
            {f'<div class="topic">{job["topic"]}</div>' if job['topic'] else ''}
            <div class="meta">{job['source'] or 'ACC/AHA'} • {job['year'] or 'N/A'} • Generated by {job['model']}</div>
        </header>
        
        <nav class="toc">
            <h2>📑 Contents</h2>
            <div>'''
    
    # Add TOC items
    for key, config in section_config.items():
        items = summary.get(key, [])
        if items:
            html += f'<a href="#{key}" class="toc-item" style="border-left: 3px solid {config["color"]};">{config["icon"]} {config["title"]} ({len(items)})</a>'
    
    html += '''
            </div>
        </nav>'''
    
    # Add sections
    total_items = 0
    for key, config in section_config.items():
        items = summary.get(key, [])
        if items:
            html += f'''
        <section class="section" id="{key}">
            <div class="section-header" onclick="this.parentElement.classList.toggle('collapsed')" style="border-left: 4px solid {config['color']};">
                <span class="section-icon">{config['icon']}</span>
                <span class="section-title">{config['title']}</span>
                <span class="section-count">{len(items)} items</span>
                <span class="section-toggle">▼</span>
            </div>
            <div class="section-content">'''
            
            for i, item in enumerate(items, 1):
                if isinstance(item, dict):
                    item_text = item.get('text', item.get('recommendation', str(item)))
                    page = item.get('page', '')
                    loe = item.get('loe', '')
                    item_class = item.get('class', '')
                else:
                    item_text = str(item)
                    page = ''
                    loe = ''
                    item_class = ''
                
                # Fix spacing issues in the text
                item_text = fix_pdf_spacing(item_text)
                
                # Escape HTML in text
                import html as html_module
                item_text = html_module.escape(item_text)
                
                html += f'''
                <div class="item" style="border-color: {config['color']};">
                    <div class="item-text"><strong>{i}.</strong> {item_text}</div>
                    <div class="item-meta">'''
                
                if page:
                    html += f'<span class="badge badge-page">Page {page}</span>'
                if item_class:
                    class_lower = item_class.lower().replace(' ', '')
                    badge_class = 'badge-class-i' if class_lower == 'i' else f'badge-class-{class_lower}'
                    html += f'<span class="badge {badge_class}">Class {item_class}</span>'
                if loe:
                    html += f'<span class="badge badge-loe">LOE: {loe}</span>'
                
                # Check for figure/table reference
                if isinstance(item, dict) and item.get('has_figure'):
                    figure_ref = item.get('figure_ref', 'Table/Figure')
                    html += f'<span class="badge" style="background: rgba(251, 191, 36, 0.2); color: #fbbf24;">📊 {figure_ref}</span>'
                
                html += '''
                    </div>'''
                
                # If has figure, add image section (expanded by default)
                if isinstance(item, dict) and item.get('has_figure') and page:
                    guideline_id = job['guideline_id']
                    figure_ref = item.get('figure_ref', 'Figure')
                    html += f'''
                    <div class="figure-container" style="margin-top: 15px; text-align: center; padding: 15px; background: #0f172a; border-radius: 8px; border: 1px solid #fbbf24;">
                        <p style="margin-bottom: 10px; font-size: 0.9rem; color: #fbbf24; font-weight: 600;">📊 {figure_ref} (Page {page})</p>
                        <img src="https://cardio-guidelines.onrender.com/api/page-image/{guideline_id}/{page}" 
                             style="max-width: 100%; border-radius: 8px; box-shadow: 0 4px 20px rgba(0,0,0,0.5);"
                             alt="{figure_ref} from Page {page}"
                             loading="lazy" />
                    </div>'''
                
                html += '''
                </div>'''
                total_items += 1
            
            html += '''
            </div>
        </section>'''
    
    html += f'''
        <footer class="footer">
            <p>Total items extracted: {total_items}</p>
            <p>Generated with {job['model']} • Cost: ${job['cost']:.4f} • Tokens: {job['tokens']:,}</p>
            <p>Cardiology Guidelines Assistant</p>
        </footer>
    </div>
</body>
</html>'''
    
    # Save HTML to cache for future downloads
    cur.execute("UPDATE pptx_jobs SET html_data = %s WHERE id = %s", (html, job_id))
    conn.commit()
    print(f"Cached HTML for job {job_id} ({len(html)} chars)")
    
    cur.close()
    conn.close()
    
    return Response(
        html,
        mimetype='text/html',
        headers={'Content-Disposition': f'attachment; filename="{filename}"'}
    )

def process_pptx_job(job_id):
    """Background worker to process PPTX extraction."""
    import json
    
    conn = get_db()
    cur = conn.cursor(row_factory=dict_row)
    
    try:
        # Get job details
        cur.execute("SELECT * FROM pptx_jobs WHERE id = %s", (job_id,))
        job = cur.fetchone()
        
        if not job:
            return
        
        guideline_id = job['guideline_id']
        topic = job['topic']
        model = job['model']
        
        # Get guideline
        cur.execute("SELECT * FROM guidelines WHERE id = %s", (guideline_id,))
        guideline = cur.fetchone()
        
        # Get chunks - use page ranges from topics table if available
        if topic:
            # First check if we have saved page ranges for this topic
            cur.execute("""
                SELECT start_page, end_page, chunk_ids 
                FROM guideline_topics 
                WHERE guideline_id = %s AND topic_name = %s
            """, (guideline_id, topic))
            topic_info = cur.fetchone()
            
            if topic_info and topic_info['start_page'] and topic_info['end_page']:
                # Use page range from topic analysis
                print(f"Using saved topic pages: {topic_info['start_page']}-{topic_info['end_page']}")
                cur.execute("""
                    SELECT content, page_number FROM chunks 
                    WHERE guideline_id = %s 
                    AND page_number >= %s AND page_number <= %s
                    AND (is_reference IS NULL OR is_reference = FALSE)
                    ORDER BY page_number
                """, (guideline_id, topic_info['start_page'], topic_info['end_page']))
            else:
                # Fallback to keyword search
                print(f"No saved pages, using keyword search for: {topic}")
                cur.execute("""
                    SELECT content, page_number FROM chunks 
                    WHERE guideline_id = %s 
                    AND LOWER(content) LIKE %s
                    AND (is_reference IS NULL OR is_reference = FALSE)
                    ORDER BY page_number
                """, (guideline_id, f'%{topic.lower()}%'))
        else:
            cur.execute("""
                SELECT content, page_number FROM chunks 
                WHERE guideline_id = %s
                AND (is_reference IS NULL OR is_reference = FALSE)
                ORDER BY page_number
            """, (guideline_id,))
        
        all_chunks = cur.fetchall()
        
        if not all_chunks:
            cur.execute("UPDATE pptx_jobs SET status = 'failed', error = 'No content found' WHERE id = %s", (job_id,))
            conn.commit()
            return
        
        print(f"=== PPTX JOB {job_id} STARTED ===")
        print(f"Guideline: {guideline['name']}")
        print(f"Total chunks: {len(all_chunks)}")
        
        # Process in batches - smaller batches = faster API responses
        CHUNKS_PER_BATCH = 25
        all_content = {
            'recommendations': [],
            'definitions': [],
            'epidemiology': [],
            'pathophysiology': [],
            'diagnosis': [],
            'treatment': [],
            'management': [],
            'special_populations': [],
            'key_points': []
        }
        total_input_tokens = 0
        total_output_tokens = 0
        batch_num = 0
        topic_text = f" focusing on {topic}" if topic else ""
        
        for i in range(0, len(all_chunks), CHUNKS_PER_BATCH):
            batch_chunks = all_chunks[i:i + CHUNKS_PER_BATCH]
            batch_num += 1
            
            # Update progress
            cur.execute("UPDATE pptx_jobs SET progress = %s WHERE id = %s", (batch_num, job_id))
            conn.commit()
            
            # Apply spacing fix to chunk content
            context = "\n\n---\n\n".join([
                f"[Page {c['page_number']}]\n{fix_pdf_spacing(c['content'])}" 
                for c in batch_chunks
            ])
            
            print(f"Job {job_id} - Batch {batch_num}: chunks {i+1}-{i+len(batch_chunks)}")
            
            extraction_prompt = f"""You are extracting and organizing ALL clinical content from a medical guideline.

THIS IS NOT A SUMMARY. Extract ALL information COMPLETELY and EXACTLY as written.
DO NOT shorten, paraphrase, or omit ANY clinical content.

GUIDELINE: {guideline['name']}
PAGES: {batch_chunks[0]['page_number']}-{batch_chunks[-1]['page_number']}
{f"TOPIC FOCUS: {topic}" if topic else ""}

CONTENT TO EXTRACT FROM:
{context}

YOUR TASK: Extract and organize ALL the clinical content into these categories:

1. **recommendations** - Any Class I, IIa, IIb, III recommendations (with LOE if present)
2. **definitions** - Definitions, diagnostic criteria, classifications, staging
3. **epidemiology** - Statistics, prevalence, incidence, risk factors, prognosis
4. **pathophysiology** - Mechanisms, causes, disease processes
5. **diagnosis** - Diagnostic approaches, tests, imaging, evaluation criteria
6. **treatment** - Medical therapy, medications, doses, procedures, interventions
7. **management** - Follow-up, monitoring, surveillance, timing
8. **special_populations** - Elderly, pregnancy, pediatric, comorbidities
9. **key_points** - Important clinical pearls, warnings, practical tips

CRITICAL RULES:
1. Extract COMPLETE text - do NOT shorten or summarize
2. Include page numbers for each item
3. If content doesn't fit a category, put it in "key_points"
4. Include ALL tables, criteria, algorithms as detailed text
5. IMPORTANT: If content references a Table, Figure, or Algorithm, set "has_figure": true and "figure_ref": "Table X" or "Figure Y"

EXCLUDE THE FOLLOWING (DO NOT EXTRACT):
- Journal article references like "Smith J, et al. Circulation. 2020;142:e123-145"
- Numbered bibliography entries like "1. Author AA. Journal Name. Year;Vol:Pages"
- DOI links like "doi:10.1016/j.jacc.2020.01.001"
- Any text that looks like: "Author Name(s). Journal Abbreviation. Year;Volume:Pages"
- Inline citations like "(1)", "(15)", "[23]", "(Smith 2020)"
- Author lists and affiliations
- Acknowledgments, disclosures, conflicts of interest
- Supplementary material references

Return ONLY valid JSON:
{{
    "recommendations": [{{"text": "full text", "class": "I/IIa/IIb/III", "loe": "A/B/C", "page": 15, "has_figure": false}}],
    "definitions": [{{"text": "full text", "page": 12, "has_figure": true, "figure_ref": "Table 1"}}],
    "epidemiology": [{{"text": "full text", "page": 10, "has_figure": false}}],
    "pathophysiology": [{{"text": "full text", "page": 11, "has_figure": false}}],
    "diagnosis": [{{"text": "full text", "page": 14, "has_figure": true, "figure_ref": "Figure 3"}}],
    "treatment": [{{"text": "full text", "page": 16, "has_figure": false}}],
    "management": [{{"text": "full text", "page": 18, "has_figure": false}}],
    "special_populations": [{{"text": "full text", "page": 20, "has_figure": false}}],
    "key_points": [{{"text": "full text", "page": 22, "has_figure": false}}]
}}"""

            try:
                batch_text = None  # Initialize
                
                if model.startswith('gpt-'):
                    import requests as req
                    openai_key = os.environ.get('OPENAI_API_KEY')
                    request_body = {
                        'model': model,
                        'messages': [{'role': 'user', 'content': extraction_prompt}],
                        'temperature': 0,
                        'max_completion_tokens': 8000  # GPT-5 requires this instead of max_tokens
                    }
                    print(f"  Calling GPT API (timeout=180s)...")
                    resp = req.post(
                        'https://api.openai.com/v1/chat/completions',
                        headers={'Authorization': f'Bearer {openai_key}', 'Content-Type': 'application/json'},
                        json=request_body,
                        timeout=180  # Increased from 120
                    )
                    print(f"  GPT API response: {resp.status_code}")
                    if resp.status_code == 200:
                        resp_data = resp.json()
                        batch_text = resp_data['choices'][0]['message']['content']
                        usage = resp_data.get('usage', {})
                        total_input_tokens += usage.get('prompt_tokens', 0)
                        total_output_tokens += usage.get('completion_tokens', 0)
                    else:
                        print(f"Batch {batch_num} OpenAI error: {resp.status_code} - {resp.text[:500]}")
                        continue
                        continue
                elif model.endswith('-thinking'):
                    base_model = model.replace('-thinking', '')
                    response = claude.messages.create(
                        model=base_model,
                        max_tokens=16000,
                        thinking={"type": "enabled", "budget_tokens": 8000},
                        messages=[{"role": "user", "content": extraction_prompt}]
                    )
                    batch_text = ""
                    for block in response.content:
                        if block.type == "text":
                            batch_text = block.text
                            break
                    total_input_tokens += response.usage.input_tokens
                    total_output_tokens += response.usage.output_tokens
                else:
                    response = claude.messages.create(
                        model=model,
                        max_tokens=8000,
                        temperature=0,
                        messages=[{"role": "user", "content": extraction_prompt}]
                    )
                    batch_text = response.content[0].text
                    total_input_tokens += response.usage.input_tokens
                    total_output_tokens += response.usage.output_tokens
                
                # Check if we got a response
                if not batch_text:
                    print(f"  No response text received from API")
                    continue
                    
                # Log response length for debugging
                print(f"  API response length: {len(batch_text)} chars")
                
                # Parse results
                start = batch_text.find('{')
                end = batch_text.rfind('}') + 1
                if start >= 0 and end > start:
                    try:
                        batch_json = json.loads(batch_text[start:end])
                        items_found = 0
                        for content_key in all_content.keys():
                            batch_items = batch_json.get(content_key, [])
                            if batch_items:
                                all_content[content_key].extend(batch_items)
                                items_found += len(batch_items)
                                print(f"  Found {len(batch_items)} {content_key}")
                        if items_found == 0:
                            print(f"  WARNING: JSON parsed but no items found. Keys in response: {list(batch_json.keys())}")
                    except json.JSONDecodeError as je:
                        print(f"  JSON parse error: {je}")
                        print(f"  Raw text (first 500 chars): {batch_text[:500]}")
                else:
                    print(f"  No JSON found in response. Raw text (first 500 chars): {batch_text[:500]}")
                            
            except Exception as e:
                print(f"Batch {batch_num} error: {e}")
                import traceback
                print(traceback.format_exc())
                continue
        
        # Calculate final stats
        total_items = sum(len(v) for v in all_content.values())
        cost = calculate_cost(model.replace('-thinking', ''), total_input_tokens, total_output_tokens)
        
        print(f"=== JOB {job_id} COMPLETE ===")
        print(f"Total content items: {total_items}")
        print(f"Cost: ${cost:.4f}")
        
        if total_items == 0:
            error_msg = "No content extracted. The AI may have returned an unexpected format. Check server logs for details."
            print(f"ERROR: {error_msg}")
            cur.execute("""
                UPDATE pptx_jobs 
                SET status = 'failed', error = %s
                WHERE id = %s
            """, (error_msg, job_id))
        else:
            summary_json = {
                'title': guideline['name'],
                'subtitle': topic or 'Complete Content',
                **all_content
            }
            
            cur.execute("""
                UPDATE pptx_jobs 
                SET status = 'completed', result_data = %s, cost = %s, tokens = %s, completed_at = NOW()
                WHERE id = %s
            """, (json.dumps(summary_json), cost, total_input_tokens + total_output_tokens, job_id))
        
        conn.commit()
        
    except Exception as e:
        print(f"Job {job_id} failed: {e}")
        import traceback
        print(traceback.format_exc())
        cur.execute("UPDATE pptx_jobs SET status = 'failed', error = %s WHERE id = %s", (str(e), job_id))
        conn.commit()
    finally:
        cur.close()
        conn.close()


def create_pptx(summary, guideline, topic=None):
    """Create a beautiful PowerPoint file from comprehensive content data."""
    if not PPTX_AVAILABLE:
        print("PPTX generation skipped - python-pptx not available")
        return None
    
    try:
        from pptx.enum.shapes import MSO_SHAPE
        import io
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Color palette
        NAVY = RGBColor(15, 23, 42)
        WHITE = RGBColor(255, 255, 255)
        LIGHT_GRAY = RGBColor(226, 232, 240)
        
        # Section colors
        SECTION_COLORS = {
            'recommendations': RGBColor(34, 197, 94),    # Green
            'definitions': RGBColor(59, 130, 246),       # Blue
            'epidemiology': RGBColor(168, 85, 247),      # Purple
            'pathophysiology': RGBColor(236, 72, 153),   # Pink
            'diagnosis': RGBColor(251, 191, 36),         # Amber
            'treatment': RGBColor(20, 184, 166),         # Teal
            'management': RGBColor(249, 115, 22),        # Orange
            'special_populations': RGBColor(139, 92, 246), # Violet
            'key_points': RGBColor(239, 68, 68),         # Red
        }
        
        SECTION_TITLES = {
            'recommendations': 'Clinical Recommendations',
            'definitions': 'Definitions & Classifications',
            'epidemiology': 'Epidemiology & Risk Factors',
            'pathophysiology': 'Pathophysiology',
            'diagnosis': 'Diagnosis & Evaluation',
            'treatment': 'Treatment & Therapy',
            'management': 'Management & Follow-up',
            'special_populations': 'Special Populations',
            'key_points': 'Key Clinical Points',
        }
        
        # Helper function to get page image
        def get_page_image_bytes(page_number):
            """Get page image as bytes for embedding in PPTX."""
            try:
                conn = get_db()
                cur = conn.cursor(row_factory=dict_row)
                cur.execute("SELECT pdf_data FROM guidelines WHERE id = %s", (guideline['id'],))
                result = cur.fetchone()
                cur.close()
                conn.close()
                
                if not result or not result['pdf_data']:
                    print(f"  No PDF data found for guideline {guideline['id']}")
                    return None
                
                pdf_bytes = bytes(result['pdf_data'])
                
                # Try pdf2image first
                try:
                    from pdf2image import convert_from_bytes
                    print(f"  Converting page {page_number} with pdf2image...")
                    images = convert_from_bytes(
                        pdf_bytes,
                        first_page=page_number,
                        last_page=page_number,
                        dpi=150
                    )
                    if images:
                        img_buffer = io.BytesIO()
                        images[0].save(img_buffer, format='PNG')
                        img_buffer.seek(0)
                        print(f"  Successfully converted page {page_number} with pdf2image")
                        return img_buffer
                except ImportError:
                    print(f"  pdf2image not available, trying PyMuPDF...")
                except Exception as e:
                    print(f"  pdf2image error: {e}, trying PyMuPDF...")
                
                # Fallback to PyMuPDF
                try:
                    import fitz  # PyMuPDF
                    print(f"  Converting page {page_number} with PyMuPDF...")
                    pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                    page = pdf_doc[page_number - 1]  # 0-indexed
                    
                    # Render at 2x for good quality
                    mat = fitz.Matrix(2, 2)
                    pix = page.get_pixmap(matrix=mat)
                    
                    img_buffer = io.BytesIO(pix.tobytes("png"))
                    img_buffer.seek(0)
                    pdf_doc.close()
                    print(f"  Successfully converted page {page_number} with PyMuPDF")
                    return img_buffer
                except Exception as mupdf_err:
                    print(f"  PyMuPDF error: {mupdf_err}")
                    return None
                    
            except Exception as db_err:
                print(f"  Database error getting PDF: {db_err}")
                return None
        
        # Track pages we've already added images for (avoid duplicates)
        pages_with_images = set()
        
        # =====================
        # TITLE SLIDE
        # =====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = NAVY
        bg_shape.line.fill.background()
        
        top_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        top_stripe.fill.solid()
        top_stripe.fill.fore_color.rgb = SECTION_COLORS['recommendations']
        top_stripe.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.733), Inches(1.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = summary.get('title', guideline['name'])
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.733), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = summary.get('subtitle', topic or 'Complete Content')
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = LIGHT_GRAY
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        source_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.5))
        source_frame = source_box.text_frame
        source_para = source_frame.paragraphs[0]
        source_para.text = f"{guideline.get('source', 'ACC/AHA')} {guideline.get('year', '')} Guidelines"
        source_para.font.size = Pt(16)
        source_para.font.color.rgb = RGBColor(148, 163, 184)
        source_para.alignment = PP_ALIGN.CENTER
        
        # =====================
        # TABLE OF CONTENTS
        # =====================
        toc_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = toc_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
        
        toc_title = toc_slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
        toc_title_frame = toc_title.text_frame
        toc_title_para = toc_title_frame.paragraphs[0]
        toc_title_para.text = "Contents"
        toc_title_para.font.size = Pt(36)
        toc_title_para.font.bold = True
        toc_title_para.font.color.rgb = WHITE
        
        toc_top = 1.5
        for section_key, section_title in SECTION_TITLES.items():
            items = summary.get(section_key, [])
            if items:
                toc_item = toc_slide.shapes.add_textbox(Inches(1.0), Inches(toc_top), Inches(10), Inches(0.5))
                toc_frame = toc_item.text_frame
                toc_para = toc_frame.paragraphs[0]
                toc_para.text = f"{section_title} ({len(items)} items)"
                toc_para.font.size = Pt(20)
                toc_para.font.color.rgb = SECTION_COLORS.get(section_key, WHITE)
                toc_top += 0.6
        
        # =====================
        # CONTENT SLIDES
        # =====================
        for section_key, section_title in SECTION_TITLES.items():
            items = summary.get(section_key, [])
            if not items:
                continue
            
            section_color = SECTION_COLORS.get(section_key, WHITE)
            items_per_slide = 2
            
            for slide_num, start_idx in enumerate(range(0, len(items), items_per_slide)):
                slide_items = items[start_idx:start_idx + items_per_slide]
                total_slides = -(-len(items) // items_per_slide)
                
                content_slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                bg = content_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
                bg.fill.solid()
                bg.fill.fore_color.rgb = NAVY
                bg.line.fill.background()
                
                top_bar = content_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = section_color
                top_bar.line.fill.background()
                
                left_bar = content_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.5), Inches(0.12), Inches(0.8))
                left_bar.fill.solid()
                left_bar.fill.fore_color.rgb = section_color
                left_bar.line.fill.background()
                
                header_box = content_slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(10), Inches(0.6))
                header_frame = header_box.text_frame
                header_para = header_frame.paragraphs[0]
                header_para.text = section_title
                header_para.font.size = Pt(28)
                header_para.font.bold = True
                header_para.font.color.rgb = section_color
                
                if total_slides > 1:
                    page_box = content_slide.shapes.add_textbox(Inches(11), Inches(0.5), Inches(2), Inches(0.4))
                    page_frame = page_box.text_frame
                    page_para = page_frame.paragraphs[0]
                    page_para.text = f"{slide_num + 1}/{total_slides}"
                    page_para.font.size = Pt(14)
                    page_para.font.color.rgb = RGBColor(100, 116, 139)
                    page_para.alignment = PP_ALIGN.RIGHT
                
                current_top = 1.4
                for i, item in enumerate(slide_items):
                    if isinstance(item, dict):
                        item_text = item.get('text', item.get('recommendation', str(item)))
                        page_num = item.get('page', '')
                        loe = item.get('loe', '')
                        item_class = item.get('class', '')
                    else:
                        item_text = str(item)
                        page_num = ''
                        loe = ''
                        item_class = ''
                    
                    # Fix spacing issues in the text
                    item_text = fix_pdf_spacing(item_text)
                    
                    num_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(current_top), Inches(0.5), Inches(0.4))
                    num_frame = num_box.text_frame
                    num_para = num_frame.paragraphs[0]
                    num_para.text = f"{start_idx + i + 1}."
                    num_para.font.size = Pt(16)
                    num_para.font.bold = True
                    num_para.font.color.rgb = section_color
                    
                    # Main text with formatting
                    text_box = content_slide.shapes.add_textbox(Inches(1.0), Inches(current_top), Inches(11.5), Inches(2.2))
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    text_para = text_frame.paragraphs[0]
                    
                    # Truncate if needed
                    display_text = item_text[:800] if len(item_text) > 800 else item_text
                    
                    # Add formatted text with runs for highlighting key terms
                    import re
                    
                    # Find patterns to highlight
                    # Pattern for drug names (common cardiac drugs)
                    drug_pattern = r'\b(aspirin|warfarin|apixaban|rivaroxaban|dabigatran|edoxaban|clopidogrel|ticagrelor|prasugrel|metoprolol|carvedilol|bisoprolol|lisinopril|enalapril|ramipril|losartan|valsartan|sacubitril|amlodipine|diltiazem|verapamil|amiodarone|sotalol|dofetilide|flecainide|propafenone|digoxin|furosemide|spironolactone|eplerenone|atorvastatin|rosuvastatin|heparin|enoxaparin)\b'
                    
                    # Split text and add runs with formatting
                    text_para.text = ""  # Clear default
                    
                    # Simple approach: add the main text, then badges below
                    run = text_para.add_run()
                    run.text = display_text
                    run.font.size = Pt(15)
                    run.font.color.rgb = WHITE
                    
                    text_para.line_spacing = 1.3
                    
                    # Add Class/LOE badges as colored boxes below the text
                    badge_top = current_top + 2.0
                    badge_left = 1.0
                    
                    # Class badge
                    if item_class:
                        class_colors = {
                            'i': RGBColor(34, 197, 94),      # Green
                            'I': RGBColor(34, 197, 94),
                            'iia': RGBColor(59, 130, 246),   # Blue
                            'IIa': RGBColor(59, 130, 246),
                            'iib': RGBColor(251, 191, 36),   # Yellow/Amber
                            'IIb': RGBColor(251, 191, 36),
                            'iii': RGBColor(239, 68, 68),    # Red
                            'III': RGBColor(239, 68, 68),
                        }
                        class_color = class_colors.get(item_class, RGBColor(148, 163, 184))
                        
                        class_box = content_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(badge_left), Inches(badge_top), Inches(1.2), Inches(0.35))
                        class_box.fill.solid()
                        class_box.fill.fore_color.rgb = class_color
                        class_box.line.fill.background()
                        
                        class_tf = class_box.text_frame
                        class_tf.paragraphs[0].text = f"Class {item_class}"
                        class_tf.paragraphs[0].font.size = Pt(11)
                        class_tf.paragraphs[0].font.bold = True
                        class_tf.paragraphs[0].font.color.rgb = WHITE if item_class.upper() not in ['IIB'] else RGBColor(0, 0, 0)
                        class_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                        
                        badge_left += 1.4
                    
                    # LOE badge
                    if loe:
                        loe_box = content_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(badge_left), Inches(badge_top), Inches(1.0), Inches(0.35))
                        loe_box.fill.solid()
                        loe_box.fill.fore_color.rgb = RGBColor(139, 92, 246)  # Purple
                        loe_box.line.fill.background()
                        
                        loe_tf = loe_box.text_frame
                        loe_tf.paragraphs[0].text = f"LOE: {loe}"
                        loe_tf.paragraphs[0].font.size = Pt(11)
                        loe_tf.paragraphs[0].font.bold = True
                        loe_tf.paragraphs[0].font.color.rgb = WHITE
                        loe_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                        
                        badge_left += 1.2
                    
                    # Page number badge
                    if page_num:
                        page_box = content_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(badge_left), Inches(badge_top), Inches(1.0), Inches(0.35))
                        page_box.fill.solid()
                        page_box.fill.fore_color.rgb = RGBColor(71, 85, 105)  # Slate
                        page_box.line.fill.background()
                        
                        page_tf = page_box.text_frame
                        page_tf.paragraphs[0].text = f"p. {page_num}"
                        page_tf.paragraphs[0].font.size = Pt(11)
                        page_tf.paragraphs[0].font.color.rgb = WHITE
                        page_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    current_top += 2.8
                    
                    # Check if this item has a figure/table reference - add page image slide
                    if isinstance(item, dict) and item.get('has_figure') and page_num:
                        page_int = int(page_num) if str(page_num).isdigit() else None
                        if page_int and page_int not in pages_with_images:
                            img_bytes = get_page_image_bytes(page_int)
                            if img_bytes:
                                pages_with_images.add(page_int)
                                figure_ref = item.get('figure_ref', f'Page {page_int}')
                                
                                # Create figure slide
                                fig_slide = prs.slides.add_slide(prs.slide_layouts[6])
                                
                                fig_bg = fig_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
                                fig_bg.fill.solid()
                                fig_bg.fill.fore_color.rgb = NAVY
                                fig_bg.line.fill.background()
                                
                                # Title bar
                                fig_bar = fig_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
                                fig_bar.fill.solid()
                                fig_bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
                                fig_bar.line.fill.background()
                                
                                fig_title = fig_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
                                fig_title_frame = fig_title.text_frame
                                fig_title_para = fig_title_frame.paragraphs[0]
                                fig_title_para.text = f"📊 {figure_ref} (Page {page_int})"
                                fig_title_para.font.size = Pt(24)
                                fig_title_para.font.bold = True
                                fig_title_para.font.color.rgb = section_color
                                
                                # Add the image
                                try:
                                    fig_slide.shapes.add_picture(img_bytes, Inches(0.5), Inches(1.0), width=Inches(12.333), height=Inches(6.2))
                                    print(f"  Added page image for {figure_ref} (Page {page_int})")
                                except Exception as img_err:
                                    print(f"  Failed to add image: {img_err}")
        
        output_path = f'/tmp/pptx_{guideline["id"]}_{topic or "full"}.pptx'
        prs.save(output_path)
        return output_path
        
    except Exception as e:
        print(f"PPTX creation error: {e}")
        import traceback
        print(traceback.format_exc())
        return None


# =============================================================================
# MAIN
# =============================================================================

# Initialize database on startup
try:
    init_db()
    print("Database initialized successfully")
except Exception as e:
    print(f"Database initialization error: {e}")

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=True)
